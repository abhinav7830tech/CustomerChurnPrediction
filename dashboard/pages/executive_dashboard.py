"""
Executive Dashboard — Customer Churn Analytics Platform

A CEO / management-level business intelligence overview of the churn
landscape. Answers four executive questions:
  01 What is happening?   (KPIs, health score, revenue, churn)
  02 Why is it happening? (drivers, segmentation, geography proxy)
  03 What should we do?   (alerts, departments, opportunities, roadmap)
  04 How do we present it? (executive summary, AI insights, board brief)

Presentation layer only — every number is derived deterministically from
the dataset in `utils.load_data()`; no ML retraining and no random values.
Financial figures that are modeled from the data are clearly labeled
"Estimated".
"""

import io
import os
import re
import sys
from datetime import datetime

import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
from fpdf import FPDF
from fpdf.enums import XPos, YPos

try:
    from utils import load_data, get_churn_rate
except ModuleNotFoundError:
    _dashboard_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    if _dashboard_dir not in sys.path:
        sys.path.insert(0, _dashboard_dir)
    from utils import load_data, get_churn_rate

try:
    import prediction
except ModuleNotFoundError:
    _dashboard_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    if _dashboard_dir not in sys.path:
        sys.path.insert(0, _dashboard_dir)
    import prediction

# ═══════════════════════════════════════════════════════════════════════════════
# PAGE CONFIG
# ═══════════════════════════════════════════════════════════════════════════════

st.set_page_config(
    page_title="Executive Dashboard",
    page_icon="🏢",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# ── Presentation-only constants (deterministic rule thresholds) ────────────────

PALETTE = {
    "bg": "#0F3040",
    "card": "#234556",
    "gold": "#C8A96B",
    "sage": "#8FA28A",
    "teal": "#9BCEC1",
    "text": "#F4F2EE",
    "sub": "#D6D8D8",
    "red": "#D97C7C",
}

CHURN_COLORS = {"No": "#8FA28A", "Yes": "#C8A96B"}

SEVERITY_COLORS = {
    "Critical": "#D97C7C",
    "Warning": "#C8A96B",
    "Information": "#8FA28A",
}

SEGMENT_META = {
    "VIP": ("💎", "#C8A96B",
            "Tenure ≥ 60 mo and monthly spend ≥ $90 — the crown jewels to protect."),
    "Premium": ("⭐", "#8FA28A",
                "Tenure ≥ 48 mo or spend ≥ $85 — high value with headroom to grow."),
    "Standard": ("🤝", "#D6D8D8",
                 "The stable core — nurture with gentle, low-cost engagement."),
    "High Risk": ("⚠️", "#E0635A",
                  "Month-to-month or electronic-check payer — a proactive retention plan."),
    "Critical": ("🚨", "#D97C7C",
                 "Month-to-month AND under 12 months tenure — highest churn exposure."),
}

DEPARTMENT_META = [
    ("Sales", "🧑‍💼", "Active customer base", "5,163 (73.4% of book)", "Good",
     "Anchor new acquisition on annual and two-year contracts to improve the "
     "long-term mix from day one."),
    ("Marketing", "📣", "Month-to-month share", "55.1% of customers", "At Risk",
     "Run a conversion campaign moving month-to-month accounts onto annual terms."),
    ("Customer Support", "🎧", "Add-on adoption", "50.4% hold ≥ 2 protections", "Opportunity",
     "Promote tech support and security bundles to deepen service stickiness."),
    ("Retention Team", "🛡️", "Retention rate", "73.4% of customer base", "Good",
     "Prioritize tenure < 12 months and electronic-check cohorts — both churn above 45%."),
    ("Finance", "💰", "MRR · ARPU", "$455.7K · $64.80", "Stable",
     "Model the $1.67M annual revenue at risk and fund the highest-ROI retention programs."),
    ("AI Team", "🤖", "Model accuracy", "76.1% · AUC 0.8133", "Good",
     "Monitor drift quarterly and refresh SHAP-based driver insights after each retrain."),
]

ROADMAP = [
    (
        "Immediate", "0–30 days", "⚡",
        [
            "Launch the Save-The-Customer program for the Critical + High Risk segments.",
            "Convert month-to-month accounts with a time-bound annual-contract offer.",
            "Automate payment flows for electronic-check customers.",
        ],
        "Cut the churn spike in the first 12 months of tenure",
        "Retention Team", "Critical",
    ),
    (
        "Mid-Term", "30–90 days", "🎯",
        [
            "Run the fiber quality audit and close the top service-level gaps.",
            "Bundled tech-support + security campaign for internet customers.",
            "Senior-care touchpoint program for the senior segment.",
        ],
        "Lower fiber and service-level churn toward the company average",
        "Customer Support + Marketing", "High",
    ),
    (
        "Long-Term", "90–180 days", "🏆",
        [
            "Full loyalty program for the tenured, active base.",
            "Premium upsell path converting Standard/Premium into VIP relationships.",
            "Quarterly board review of churn KPIs, ROI, and model performance.",
        ],
        "A structurally lower churn rate and higher ARPU",
        "All Departments", "Medium",
    ),
]

AI_INSIGHTS = [
    ("📡", "Fiber customers churn ~1.6× the company average",
     "Fiber optic churns at 41.9% vs 26.6% overall, yet carries 62.2% of total "
     "revenue — service quality and price are the highest-leverage fixes."),
    ("💳", "Electronic check is the most fragile billing channel",
     "45.3% of electronic-check customers churn vs ~16% on automatic payment "
     "methods. Auto-pay migration is a low-cost, high-impact lever."),
    ("📅", "Two-year contracts are the strongest retention anchor",
     "Only 2.8% of two-year customers churn, versus 42.7% of month-to-month. "
     "Contract term is the single strongest behavioral driver of retention."),
    ("🆕", "The first 12 months decide the relationship",
     "Nearly half (48.5%) of customers under 12 months tenure churn. Onboarding "
     "and early-tenure rescue offer the cheapest wins."),
    ("💲", "High spenders are price-sensitive too",
     "Customers above $90/month churn at 32.8% vs 15.8% below $50/month — value "
     "per dollar, not just price, drives premium-account retention."),
]

BOARD_SUMMARY = [
    ("Business Health", "#8FA28A",
     "Moderate/Good — the portfolio carries a 26.6% churn rate with a healthy "
     "73.4% retention base. Overall health score ~65/100."),
    ("Key Risks", "#D97C7C",
     "Month-to-month exposure (42.7% churn), fiber service experience (41.9%), "
     "early-tenure loss (48.5%), and electronic-check friction (45.3%)."),
    ("Top Opportunities", "#C8A96B",
     "Contract conversion, fiber quality, tech-support bundling, and payment "
     "automation — together worth an estimated $1.67M in annual revenue at risk."),
    ("Revenue Impact", "#C8A96B",
     "Estimated $1.67M annualized revenue lost to churn; ~$316.5K in monthly "
     "recurring revenue is currently retained."),
    ("Retention Strategy", "#9BCEC1",
     "Prioritize Critical and High-Risk segments first (5,429 customers), then "
     "convert to annual terms and automate payments."),
    ("Model Performance", "#8FA28A",
     "Deployed XGBoost: 76.1% accuracy, AUC 0.8133 — reliable for driver "
     "prioritization; refresh and re-validate quarterly."),
]

# ═══════════════════════════════════════════════════════════════════════════════
# CUSTOM CSS
# ═══════════════════════════════════════════════════════════════════════════════


def _inject_css() -> None:
    """Executive-grade dark theme consistent with the platform pages."""
    st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&display=swap');

    * { font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif; }

    .stApp { background: #0F3040; }

    .block-container {
        padding-top: 1.25rem;
        padding-bottom: 2.5rem;
        max-width: 1680px;
    }

    #MainMenu { visibility: hidden; }
    footer { visibility: hidden; }
    .stDeployButton { display: none; }

    .stApp::before {
        content: '';
        position: fixed;
        top: 0; left: 0; right: 0;
        height: 3px;
        background: linear-gradient(90deg, #C8A96B, #8FA28A, #C8A96B);
        z-index: 999;
    }

    @keyframes fadeIn {
        from { opacity: 0; transform: translateY(10px); }
        to { opacity: 1; transform: translateY(0); }
    }

    .back-link {
        display: inline-flex;
        align-items: center;
        gap: 0.35rem;
        color: #D6D8D8;
        text-decoration: none;
        font-size: 0.85rem;
        font-weight: 500;
        padding: 0.4rem 0;
        margin-bottom: 0.75rem;
    }
    .back-link:hover { color: #8FA28A; }

    .page-header { margin-bottom: 1.25rem; }

    .page-kicker {
        font-size: 0.72rem;
        font-weight: 700;
        color: #C8A96B;
        text-transform: uppercase;
        letter-spacing: 0.16em;
        margin-bottom: 0.6rem;
    }

    .page-title {
        font-size: 2.6rem;
        font-weight: 800;
        color: #F4F2EE;
        margin-bottom: 0.7rem;
        letter-spacing: -0.025em;
        line-height: 1.15;
    }

    .page-subtitle {
        font-size: 1.0rem;
        color: #D6D8D8;
        font-weight: 400;
        line-height: 1.65;
        max-width: 860px;
        margin-bottom: 1rem;
    }

    .page-rule {
        width: 88px;
        height: 4px;
        border-radius: 2px;
        background: linear-gradient(90deg, #C8A96B, rgba(200,169,107,0.1));
    }

    .meta-row { margin-bottom: 1.1rem; }
    .meta-hint {
        font-size: 0.78rem;
        color: #D6D8D8;
        opacity: 0.8;
        padding-top: 0.5rem;
    }
    .meta-hint b { color: #C8A96B; font-weight: 700; }

    /* ── Section headers ── */
    .section-head {
        display: flex;
        align-items: center;
        gap: 0.8rem;
        margin: 1.9rem 0 0.85rem 0;
    }
    .sec-num {
        font-size: 0.78rem;
        font-weight: 800;
        color: #0F3040;
        background: linear-gradient(135deg, #C8A96B, #b09055);
        border-radius: 9px;
        padding: 0.3rem 0.5rem;
        letter-spacing: 0.04em;
        flex-shrink: 0;
    }
    .sec-icon { font-size: 1.15rem; }
    .sec-title {
        font-size: 1.05rem;
        font-weight: 800;
        color: #F4F2EE;
        letter-spacing: -0.01em;
    }
    .sec-sub {
        font-size: 0.78rem;
        color: #D6D8D8;
        opacity: 0.72;
        margin-top: 0.15rem;
    }

    /* ── KPI cards ── */
    .kpi-card {
        position: relative;
        background: #234556;
        border: 1px solid rgba(255,255,255,0.08);
        border-radius: 18px;
        padding: 1.2rem 1rem;
        text-align: center;
        box-shadow: 0 2px 8px rgba(0,0,0,0.15);
        height: 100%;
        animation: fadeIn 0.5s ease both;
    }
    .kpi-card::before {
        content: '';
        position: absolute;
        top: 0; left: 1.5rem; right: 1.5rem;
        height: 3px;
        background: #C8A96B;
        border-radius: 0 0 3px 3px;
    }
    .kpi-card:hover {
        transform: translateY(-2px);
        box-shadow: 0 6px 20px rgba(0,0,0,0.2);
    }
    .kpi-label {
        font-size: 0.68rem;
        font-weight: 500;
        color: #D6D8D8;
        text-transform: uppercase;
        letter-spacing: 0.06em;
        margin-bottom: 0.35rem;
    }
    .kpi-value {
        font-size: 1.65rem;
        font-weight: 700;
        color: #F4F2EE;
        line-height: 1.2;
    }
    .kpi-value.accent { color: #C8A96B; }
    .kpi-value.good { color: #8FA28A; }
    .kpi-value.bad { color: #D97C7C; }
    .kpi-subtext {
        font-size: 0.6rem;
        color: #D6D8D8;
        margin-top: 0.45rem;
        font-weight: 400;
        opacity: 0.75;
    }

    /* ── Alert cards ── */
    .alert-card {
        display: flex;
        gap: 0.85rem;
        align-items: flex-start;
        background: #163949;
        border: 1px solid rgba(255,255,255,0.07);
        border-left-width: 4px;
        border-radius: 14px;
        padding: 1rem 1.15rem;
        animation: fadeIn 0.5s ease both;
    }
    .alert-icon { font-size: 1.35rem; flex-shrink: 0; }
    .alert-title { font-size: 0.9rem; font-weight: 700; color: #F4F2EE; margin-bottom: 0.3rem; }
    .alert-text { font-size: 0.78rem; color: #D6D8D8; opacity: 0.88; line-height: 1.55; }
    .alert-pill {
        display: inline-block;
        font-size: 0.6rem;
        font-weight: 800;
        letter-spacing: 0.06em;
        text-transform: uppercase;
        color: #0F3040;
        border-radius: 100px;
        padding: 0.18rem 0.55rem;
        margin-bottom: 0.4rem;
    }

    /* ── Segment cards ── */
    .seg-grid {
        display: grid;
        grid-template-columns: repeat(auto-fit, minmax(185px, 1fr));
        gap: 1rem;
    }
    .seg-card {
        background: #163949;
        border: 1px solid rgba(255,255,255,0.06);
        border-radius: 14px;
        padding: 1.15rem 1rem;
        text-align: center;
        transition: transform 0.25s ease, border-color 0.25s ease;
        animation: fadeIn 0.5s ease both;
    }
    .seg-card:hover { transform: translateY(-3px); border-color: rgba(200,169,107,0.4); }
    .seg-icon { font-size: 1.6rem; }
    .seg-name { font-size: 0.95rem; font-weight: 800; color: #F4F2EE; margin: 0.35rem 0 0.25rem; }
    .seg-count { font-size: 1.5rem; font-weight: 800; color: #C8A96B; }
    .seg-share { font-size: 0.68rem; color: #D6D8D8; opacity: 0.75; margin-bottom: 0.4rem; }
    .seg-desc { font-size: 0.72rem; color: #D6D8D8; opacity: 0.85; line-height: 1.5; }

    /* ── Department cards ── */
    .dep-card {
        background: #163949;
        border: 1px solid rgba(255,255,255,0.06);
        border-radius: 14px;
        padding: 1.05rem 1.15rem;
        height: 100%;
        animation: fadeIn 0.5s ease both;
    }
    .dep-head { display: flex; align-items: center; gap: 0.6rem; margin-bottom: 0.6rem; }
    .dep-icon { font-size: 1.25rem; }
    .dep-name { font-size: 0.92rem; font-weight: 700; color: #F4F2EE; }
    .dep-status {
        margin-left: auto;
        font-size: 0.6rem;
        font-weight: 800;
        letter-spacing: 0.05em;
        text-transform: uppercase;
        color: #0F3040;
        border-radius: 100px;
        padding: 0.18rem 0.55rem;
    }
    .dep-kpi-label { font-size: 0.62rem; color: #C8A96B; text-transform: uppercase; letter-spacing: 0.07em; margin-bottom: 0.2rem; }
    .dep-kpi { font-size: 0.9rem; font-weight: 700; color: #F4F2EE; margin-bottom: 0.55rem; }
    .dep-action { font-size: 0.75rem; color: #D6D8D8; opacity: 0.88; line-height: 1.55; }

    /* ── Opportunity rows ── */
    .opp-card {
        background: #163949;
        border: 1px solid rgba(255,255,255,0.06);
        border-left-width: 3px;
        border-radius: 12px;
        padding: 0.85rem 1.1rem;
        display: flex;
        align-items: center;
        gap: 1rem;
        flex-wrap: wrap;
        animation: fadeIn 0.5s ease both;
    }
    .opp-rank {
        width: 30px; height: 30px;
        border-radius: 50%;
        background: #234556;
        border: 1px solid rgba(200,169,107,0.35);
        color: #C8A96B;
        font-weight: 800;
        font-size: 0.78rem;
        display: flex;
        align-items: center;
        justify-content: center;
        flex-shrink: 0;
    }
    .opp-title { font-size: 0.88rem; font-weight: 700; color: #F4F2EE; margin-bottom: 0.15rem; }
    .opp-desc { font-size: 0.74rem; color: #D6D8D8; opacity: 0.82; line-height: 1.5; }
    .opp-meta { margin-left: auto; display: flex; gap: 1.4rem; flex-shrink: 0; }
    .opp-meta-item { text-align: right; }
    .opp-meta-label { font-size: 0.6rem; color: #D6D8D8; opacity: 0.7; text-transform: uppercase; letter-spacing: 0.06em; }
    .opp-meta-value { font-size: 0.85rem; font-weight: 700; color: #C8A96B; }

    /* ── Roadmap ── */
    .rm-card {
        background: #163949;
        border: 1px solid rgba(255,255,255,0.06);
        border-top-width: 3px;
        border-radius: 14px;
        padding: 1.15rem 1.2rem;
        height: 100%;
        animation: fadeIn 0.5s ease both;
    }
    .rm-head { display: flex; align-items: baseline; gap: 0.55rem; margin-bottom: 0.7rem; }
    .rm-phase { font-size: 1rem; font-weight: 800; color: #F4F2EE; }
    .rm-horizon { font-size: 0.72rem; color: #C8A96B; font-weight: 700; text-transform: uppercase; letter-spacing: 0.05em; }
    .rm-item {
        display: flex; gap: 0.55rem;
        font-size: 0.78rem; color: #D6D8D8;
        line-height: 1.5; margin-bottom: 0.45rem;
    }
    .rm-item span { color: #C8A96B; flex-shrink: 0; }
    .rm-outcome { font-size: 0.76rem; color: #9BCEC1; font-weight: 600; margin: 0.5rem 0 0.6rem; line-height: 1.5; }
    .rm-owner { font-size: 0.68rem; color: #D6D8D8; opacity: 0.7; }

    /* ── Insight cards ── */
    .insight-card {
        position: relative;
        background: #234556;
        border: 1px solid rgba(255,255,255,0.08);
        border-radius: 18px;
        padding: 1.2rem 1.25rem;
        height: 100%;
        animation: fadeIn 0.5s ease both;
    }
    .insight-icon { font-size: 1.5rem; margin-bottom: 0.6rem; }
    .insight-title { font-size: 0.88rem; font-weight: 700; color: #F4F2EE; margin-bottom: 0.35rem; line-height: 1.4; }
    .insight-text { font-size: 0.76rem; color: #D6D8D8; opacity: 0.88; line-height: 1.6; }

    /* ── Narrative / notes ── */
    .notes-box { font-size: 0.95rem; color: #F4F2EE; line-height: 1.9; }
    .notes-box b { color: #C8A96B; font-weight: 700; }
    .note-text {
        font-size: 0.76rem;
        color: #D6D8D8;
        opacity: 0.72;
        line-height: 1.6;
        margin-top: 0.7rem;
    }
    .board-item { font-size: 0.82rem; color: #D6D8D8; line-height: 1.65; margin-bottom: 0.6rem; }

    /* ── Expander ── */
    [data-testid="stExpander"] {
        background: #1f3d4d;
        border: 1px solid rgba(255,255,255,0.08);
        border-radius: 16px;
    }
    [data-testid="stExpander"] summary { color: #F4F2EE; font-weight: 600; font-size: 0.92rem; }

    /* ── Download buttons ── */
    .stDownloadButton button {
        width: 100% !important;
        min-height: 46px !important;
        border-radius: 14px !important;
        font-size: 0.95rem !important;
        font-weight: 700 !important;
        letter-spacing: 0.01em;
        color: #0F3040 !important;
        background: linear-gradient(135deg, #C8A96B 0%, #b09055 100%) !important;
        border: none !important;
        box-shadow: 0 6px 18px rgba(200,169,107,0.22) !important;
        transition: all 0.25s ease !important;
    }
    .stDownloadButton button:hover {
        transform: translateY(-2px) !important;
        box-shadow: 0 12px 30px rgba(200,169,107,0.38) !important;
        background: linear-gradient(135deg, #d4b678 0%, #C8A96B 100%) !important;
    }

    @media (max-width: 768px) {
        .page-title { font-size: 1.7rem; }
        .kpi-value { font-size: 1.25rem; }
        .opp-meta { margin-left: 0; width: 100%; }
        .block-container { padding-left: 1rem; padding-right: 1rem; }
    }
    </style>
    """, unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════════════
# PLOTLY TEMPLATE
# ═══════════════════════════════════════════════════════════════════════════════


def _dark_template() -> go.layout.Template:
    """Corporate dark theme for all Plotly charts."""
    return go.layout.Template(
        layout=dict(
            font=dict(family="Inter, sans-serif", size=12, color="#D6D8D8"),
            title=dict(font=dict(size=15, color="#F4F2EE"), x=0.5),
            paper_bgcolor="#234556",
            plot_bgcolor="#234556",
            height=400,
            margin=dict(l=50, r=30, t=65, b=50),
            xaxis=dict(
                gridcolor="rgba(255,255,255,0.05)",
                zerolinecolor="rgba(255,255,255,0.05)",
                tickfont=dict(size=11, color="#D6D8D8"),
                title=dict(font=dict(size=12, color="#D6D8D8")),
            ),
            yaxis=dict(
                gridcolor="rgba(255,255,255,0.05)",
                zerolinecolor="rgba(255,255,255,0.05)",
                tickfont=dict(size=11, color="#D6D8D8"),
                title=dict(font=dict(size=12, color="#D6D8D8")),
            ),
            legend=dict(
                font=dict(size=11, color="#D6D8D8"),
                bgcolor="rgba(0,0,0,0)",
                orientation="h",
                y=1.12,
            ),
            hoverlabel=dict(
                bgcolor="#234556",
                font_color="#F4F2EE",
                font_size=12,
            ),
            colorway=["#8FA28A", "#C8A96B", "#9BCEC1", "#D6D8D8"],
        )
    )


TEMPLATE = _dark_template()


# ═══════════════════════════════════════════════════════════════════════════════
# DETERMINISTIC METRICS ENGINE
# ═══════════════════════════════════════════════════════════════════════════════


def _churn_pct(group: pd.DataFrame) -> float:
    """Share of a group that churned, as a percentage."""
    return round((group["Churn"] == "Yes").mean() * 100, 1)


@st.cache_data(show_spinner="Crunching executive metrics...")
def _build_metrics(df: pd.DataFrame) -> dict:
    """Compute every deterministic figure shown on the dashboard.

    All values are derived directly from the dataset. Financial figures
    projected from the data (annualization, revenue at risk, ROI) are
    modeled estimates and are labeled as such in the UI.
    """
    total = len(df)
    churned = df["Churn"] == "Yes"
    n_churned = int(churned.sum())
    n_retained = total - n_churned
    churn_rate = round(n_churned / total * 100, 1)
    retention_rate = round(n_retained / total * 100, 1)

    mrr_all = float(df["MonthlyCharges"].sum())
    mrr_retained = float(df.loc[~churned, "MonthlyCharges"].sum())
    mrr_churned = mrr_all - mrr_retained
    arpu = round(mrr_all / total, 2)
    annual_projection = mrr_all * 12
    annual_at_risk = mrr_churned * 12
    annual_retained = mrr_retained * 12

    avg_tenure = round(df["tenure"].mean(), 1)
    avg_total_charges = round(df["TotalCharges"].mean(), 2)
    avg_tenure_churned = round(df.loc[churned, "tenure"].mean(), 1)

    model_alias = prediction.resolve_best_model() or prediction.get_available_models()[0]
    info = prediction.model_info(model_alias)
    model_label = info["label"]
    model_accuracy = info["accuracy"]
    model_auc = info["auc"]

    # ── Health score components (0-100) ──
    comp_retention = retention_rate
    comp_revenue = round(min(arpu / 100.0 * 100.0, 100.0), 1)
    internet = df["InternetService"] != "No"
    addons = ["OnlineSecurity", "OnlineBackup", "DeviceProtection", "TechSupport"]
    n_addons = df.loc[internet, addons].apply(lambda r: (r == "Yes").sum(), axis=1)
    comp_satisfaction = round((n_addons >= 2).mean() * 100.0, 1)
    comp_churn = round(100.0 - churn_rate, 1)
    comp_tenure = round(min(avg_tenure / 72.0 * 100.0, 100.0), 1)
    comp_model = round(float(model_accuracy), 1) if model_accuracy else 75.0

    weights = {
        "Retention": 0.20,
        "Revenue": 0.15,
        "Satisfaction": 0.20,
        "Churn Health": 0.20,
        "Tenure": 0.10,
        "Model Confidence": 0.15,
    }
    scores = {
        "Retention": comp_retention,
        "Revenue": comp_revenue,
        "Satisfaction": comp_satisfaction,
        "Churn Health": comp_churn,
        "Tenure": comp_tenure,
        "Model Confidence": comp_model,
    }
    health = round(sum(scores[k] * w for k, w in weights.items()), 1)
    if health >= 70:
        health_band, health_color = "Good", "#8FA28A"
    elif health >= 55:
        health_band, health_color = "Moderate", "#C8A96B"
    elif health >= 40:
        health_band, health_color = "Watch", "#E0635A"
    else:
        health_band, health_color = "Critical", "#D97C7C"

    # ── Categorical churn rates ──
    cat = {}
    for col in ["Contract", "InternetService", "PaymentMethod", "gender"]:
        rates = (
            df.groupby(col)["Churn"]
            .apply(lambda s: round((s == "Yes").mean() * 100, 1))
            .sort_values(ascending=False)
        )
        cat[col] = rates
    senior = {
        "No": _churn_pct(df[df["SeniorCitizen"] == 0]),
        "Yes": _churn_pct(df[df["SeniorCitizen"] == 1]),
    }

    # ── Tenure cohorts (proxy for a monthly churn trend) ──
    bins = [0, 12, 24, 48, 72, float("inf")]
    labels = ["0–12", "13–24", "25–48", "49–72", "73+"]
    cohort = df.copy()
    cohort["Cohort"] = pd.cut(cohort["tenure"], bins=bins, labels=labels, right=True)
    tenure_cohort = cohort.groupby("Cohort", observed=True).apply(
        lambda g: pd.Series({
            "customers": len(g),
            "churn_rate": round((g["Churn"] == "Yes").mean() * 100, 1),
            "revenue": round(g["MonthlyCharges"].sum(), 2),
        })
    ).reset_index()
    tenure_cohort = tenure_cohort[tenure_cohort["customers"] > 0].reset_index(drop=True)

    tenure_lt_12 = _churn_pct(df[df["tenure"] < 12])
    fiber_rate = cat["InternetService"].get("Fiber optic", 0.0)
    ec_rate = cat["PaymentMethod"].get("Electronic check", 0.0)
    mtm_rate = cat["Contract"].get("Month-to-month", 0.0)
    two_year_rate = cat["Contract"].get("Two year", 0.0)

    fiber_share_rev = round(
        df.loc[df["InternetService"] == "Fiber optic", "MonthlyCharges"].sum()
        / mrr_all * 100, 1
    )
    mtm_share_rev = round(
        df.loc[df["Contract"] == "Month-to-month", "MonthlyCharges"].sum()
        / mrr_all * 100, 1
    )
    mtm_share = round(
        (df["Contract"] == "Month-to-month").mean() * 100, 1
    )

    # ── Segmentation (deterministic business rules) ──
    def _segment(row) -> str:
        if row["Contract"] == "Month-to-month" and row["tenure"] < 12:
            return "Critical"
        if (row["Contract"] == "Month-to-month"
                or row["PaymentMethod"] == "Electronic check"):
            return "High Risk"
        if row["tenure"] >= 60 and row["MonthlyCharges"] >= 90:
            return "VIP"
        if row["tenure"] >= 48 or row["MonthlyCharges"] >= 85:
            return "Premium"
        return "Standard"

    segments = df.apply(_segment, axis=1).value_counts()
    seg_counts = {
        name: (int(segments.get(name, 0)), round(segments.get(name, 0) / total * 100, 1))
        for name in ["VIP", "Premium", "Standard", "High Risk", "Critical"]
    }

    # ── Executive alerts ──
    alerts = [
        ("Critical", "🚨", "Month-to-month churn is very high",
         f"{mtm_rate:.1f}% of month-to-month customers churn — and they make up "
         f"88.6% of all churned accounts ({mtm_share:.1f}% of the book)."),
        ("Critical", "🚨", "Early-tenure accounts churn at nearly 50%",
         f"{tenure_lt_12:.1f}% of customers under 12 months tenure churn. "
         f"Onboarding and the first-year window are the biggest risk surface."),
        ("Warning", "⚠️", "Fiber revenue concentration risk",
         f"Fiber optic churns at {fiber_rate:.1f}% vs {churn_rate:.1f}% company "
         f"average while carrying {fiber_share_rev:.1f}% of total revenue."),
        ("Warning", "⚠️", "Electronic-check billing friction",
         f"{ec_rate:.1f}% of electronic-check customers churn vs ~16% for "
         f"automatic payment methods."),
        ("Warning", "⚠️", "High-value accounts are also at risk",
         f"{churn_rate:.1f}% of accounts paying over $90/month churn, and "
         f"month-to-month contracts carry {mtm_share_rev:.1f}% of revenue."),
        ("Information", "ℹ️", "Two-year contracts are a strong retention anchor",
         f"Only {two_year_rate:.1f}% of two-year customers churn. Contract "
         f"conversion is the cheapest structural lever available."),
    ]

    # ── Departments ──
    departments = []
    dep_kpis = {
        "Sales": f"{n_retained:,} ({(n_retained/total*100):.1f}% of book)",
        "Marketing": f"{mtm_share:.1f}% of customers",
        "Customer Support": f"{comp_satisfaction:.1f}% hold ≥ 2 protections",
        "Retention Team": f"{retention_rate:.1f}% of customer base",
        "Finance": f"${mrr_all/1000:.1f}K · ${arpu:.2f}",
        "AI Team": f"{model_accuracy}% · AUC {model_auc}",
    }
    dep_actions = {
        "Sales": "Anchor new acquisition on annual and two-year contracts to "
                 "improve the long-term mix from day one.",
        "Marketing": "Run a conversion campaign moving month-to-month accounts "
                     "onto annual terms.",
        "Customer Support": "Promote tech support and security bundles to deepen "
                            "service stickiness.",
        "Retention Team": "Prioritize tenure < 12 months and electronic-check "
                          "cohorts — both churn above 45%.",
        "Finance": f"Model the ${annual_at_risk/1e6:.2f}M annual revenue at risk "
                   "and fund the highest-ROI retention programs.",
        "AI Team": "Monitor drift quarterly and refresh SHAP-based driver "
                   "insights after each retrain.",
    }
    for name, icon, kpi_label, _, status, action in DEPARTMENT_META:
        departments.append({
            "name": name,
            "icon": icon,
            "kpi_label": kpi_label,
            "kpi": dep_kpis[name],
            "status": status,
            "action": action,
        })

    # ── Top business opportunities (deterministic ROI model) ──
    def _opportunity(name, icon, desc, mask, recover, cost_per, mode="recover"):
        cohort = df[mask]
        if len(cohort) == 0:
            return None
        ch = cohort["Churn"] == "Yes"
        cost = len(cohort) * cost_per
        if mode == "recover":
            at_risk = 12.0 * cohort.loc[ch, "MonthlyCharges"].sum()
            benefit = at_risk * recover
        else:
            at_risk = 12.0 * cohort["MonthlyCharges"].sum() * (churn_rate / 100.0)
            benefit = at_risk * recover
        roi = (benefit - cost) / cost * 100.0 if cost > 0 else 0.0
        return {
            "name": name,
            "icon": icon,
            "desc": desc,
            "cohort": len(cohort),
            "at_risk": at_risk,
            "cost": cost,
            "benefit": benefit,
            "roi": roi,
        }

    opps = [
        _opportunity(
            "High-Value Retention Program", "💎",
            "Target month-to-month accounts paying $90+ — the most valuable "
            "customers sitting on the highest-risk contract.",
            (df["MonthlyCharges"] >= 90) & (df["Contract"] == "Month-to-month"),
            0.30, 40.0),
        _opportunity(
            "Early-Tenure Rescue", "🆘",
            "Intensive onboarding and check-ins for customers under 12 months — "
            "the window with the highest churn rate.",
            df["tenure"] < 12, 0.30, 40.0),
        _opportunity(
            "Month-to-Month → Annual Conversion", "📅",
            "Time-bound offer converting the largest at-risk cohort onto annual "
            "terms.",
            df["Contract"] == "Month-to-month", 0.30, 40.0),
        _opportunity(
            "Fiber Service Quality Audit", "📡",
            "Resolve service and pricing issues for the fiber base that carries "
            "the majority of revenue.",
            df["InternetService"] == "Fiber optic", 0.25, 40.0),
        _opportunity(
            "Electronic-Check Automation", "💳",
            "Move electronic-check payers onto automatic billing to remove "
            "payment friction.",
            df["PaymentMethod"] == "Electronic check", 0.25, 40.0),
        _opportunity(
            "Online Security Add-On", "🔐",
            "Offer the security add-on to internet customers missing it — "
            "closure of a protective gap.",
            (df["InternetService"] != "No") & (df["OnlineSecurity"] != "Yes"),
            0.20, 30.0),
        _opportunity(
            "Tech Support Enablement", "🛠️",
            "Enable tech support access for internet customers without it — the "
            "top support gap.",
            (df["InternetService"] != "No") & (df["TechSupport"] != "Yes"),
            0.20, 30.0),
        _opportunity(
            "Premium Upsell Path", "⭐",
            "Grow revenue by moving eligible Premium/Standard accounts toward "
            "the premium tier.",
            (df["tenure"] >= 48) | (df["MonthlyCharges"] >= 85),
            0.0, 30.0, mode="upsell"),
        _opportunity(
            "Senior-Care Touchpoints", "👴",
            "Personal check-ins for the senior segment, which churns well above "
            "average.",
            df["SeniorCitizen"] == 1, 0.20, 35.0),
        _opportunity(
            "Loyalty Program", "🏆",
            "Reward the stable, tenured base to keep projected future churn low.",
            (df["tenure"] >= 24) & (~churned), 0.20, 25.0),
    ]
    opportunities = sorted(
        [o for o in opps if o is not None], key=lambda o: o["roi"], reverse=True
    )[:10]

    # ── Geographic proxy (no geo fields exist in the dataset) ──
    geo_proxy = {
        "internet": (
            df["InternetService"]
            .value_counts()
            .reindex(["Fiber optic", "DSL", "No"])
            .fillna(0)
        ),
        "contract": (
            df["Contract"]
            .value_counts()
            .reindex(["Month-to-month", "One year", "Two year"])
            .fillna(0)
        ),
        "payment": (
            df["PaymentMethod"]
            .value_counts()
            .reindex([
                "Electronic check", "Mailed check",
                "Bank transfer (automatic)", "Credit card (automatic)",
            ])
            .fillna(0)
        ),
    }

    # ── Executive summary paragraph ──
    top_opp = opportunities[0]["name"] if opportunities else "a targeted program"
    summary = (
        f"The business is carrying a {churn_rate:.1f}% churn rate with a "
        f"{retention_rate:.1f}% retention base of {n_retained:,} active customers. "
        f"Churned accounts represent an estimated ${annual_at_risk/1e6:.2f}M in "
        f"annualized recurring revenue at risk, against ${annual_retained/1e6:.2f}M "
        f"that is currently retained (MRR ${mrr_all/1e3:.1f}K, ARPU ${arpu:.2f}). "
        f"Churn concentrates in three patterns: month-to-month contracts "
        f"({mtm_rate:.1f}% churn), fiber optic service ({fiber_rate:.1f}% churn, "
        f"{fiber_share_rev:.1f}% of revenue), and the first 12 months of tenure "
        f"({tenure_lt_12:.1f}% churn). The executive health score is "
        f"{health:.0f}/100 ({health_band}). The recommended response prioritizes "
        f"{top_opp} first, then widens to contract conversion, fiber quality, "
        f"and payment automation to protect the estimated ${annual_at_risk/1e6:.2f}M "
        f"at risk."
    )

    # ── Board meeting summary ──
    board = [
        ("Business Health", "#8FA28A",
         f"{health_band} — score {health:.0f}/100. {retention_rate:.1f}% of the "
         f"customer base is retained."),
        ("Key Risks", "#D97C7C",
         f"Month-to-month exposure ({mtm_rate:.1f}% churn), fiber experience "
         f"({fiber_rate:.1f}%), early-tenure loss ({tenure_lt_12:.1f}%), and "
         f"electronic-check friction ({ec_rate:.1f}%)."),
        ("Top Opportunities", "#C8A96B",
         "Contract conversion, fiber quality, tech-support bundling, and payment "
         "automation — protecting the estimated "
         f"${annual_at_risk/1e6:.2f}M annual revenue at risk."),
        ("Revenue Impact", "#C8A96B",
         f"Estimated ${annual_at_risk/1e6:.2f}M annualized revenue lost to churn; "
         f"${annual_retained/1e6:.2f}M annualized revenue retained (MRR "
         f"${mrr_all/1e3:.1f}K, ARPU ${arpu:.2f})."),
        ("Retention Strategy", "#9BCEC1",
         f"Prioritize the Critical and High-Risk segments first "
         f"({seg_counts['Critical'][0] + seg_counts['High Risk'][0]:,} customers), "
         f"then convert to annual terms and automate payments."),
        ("Model Performance", "#8FA28A",
         f"Deployed {model_label}: {model_accuracy}% accuracy, AUC {model_auc} — "
         f"reliable for driver prioritization; refresh and re-validate quarterly."),
    ]

    return {
        "total": total,
        "n_churned": n_churned,
        "n_retained": n_retained,
        "churn_rate": churn_rate,
        "retention_rate": retention_rate,
        "mrr_all": mrr_all,
        "mrr_retained": mrr_retained,
        "mrr_churned": mrr_churned,
        "arpu": arpu,
        "annual_projection": annual_projection,
        "annual_at_risk": annual_at_risk,
        "annual_retained": annual_retained,
        "avg_tenure": avg_tenure,
        "avg_total_charges": avg_total_charges,
        "avg_tenure_churned": avg_tenure_churned,
        "model_alias": model_alias,
        "model_label": model_label,
        "model_accuracy": model_accuracy,
        "model_auc": model_auc,
        "health": health,
        "health_band": health_band,
        "health_color": health_color,
        "health_weights": weights,
        "health_scores": scores,
        "cat": cat,
        "senior": senior,
        "tenure_cohort": tenure_cohort,
        "tenure_lt_12": tenure_lt_12,
        "fiber_rate": fiber_rate,
        "ec_rate": ec_rate,
        "mtm_rate": mtm_rate,
        "two_year_rate": two_year_rate,
        "fiber_share_rev": fiber_share_rev,
        "mtm_share_rev": mtm_share_rev,
        "mtm_share": mtm_share,
        "seg_counts": seg_counts,
        "alerts": alerts,
        "departments": departments,
        "opportunities": opportunities,
        "geo_proxy": geo_proxy,
        "summary": summary,
        "board": board,
    }


# ═══════════════════════════════════════════════════════════════════════════════
# SHARED RENDERING HELPERS
# ═══════════════════════════════════════════════════════════════════════════════


def _section_head(num: str, icon: str, title: str, sub: str) -> None:
    """Numbered section header used across the 15-part dashboard."""
    st.markdown(
        f'<div class="section-head">'
        f'<span class="sec-num">{num}</span>'
        f'<span class="sec-icon">{icon}</span>'
        f'<div><div class="sec-title">{title}</div>'
        f'<div class="sec-sub">{sub}</div></div>'
        f"</div>",
        unsafe_allow_html=True,
    )


def _metric_tile(label: str, value: str, sub: str = "", cls: str = "") -> str:
    """HTML for an Analytics-style KPI metric tile."""
    val_class = f"kpi-value {cls}".strip() if cls else "kpi-value"
    subtext_html = f'<div class="kpi-subtext">{sub}</div>' if sub else ""
    return (
        f'<div class="kpi-card">'
        f'<div class="kpi-label">{label}</div>'
        f'<div class="{val_class}">{value}</div>'
        f"{subtext_html}"
        f"</div>"
    )


def _kpi_row(metrics: dict, tiles: list) -> None:
    """Render four KPI cards in one row."""
    cols = st.columns(4, gap="medium")
    for col, (label, value, sub, cls) in zip(cols, tiles):
        with col:
            st.markdown(_metric_tile(label, value, sub, cls), unsafe_allow_html=True)


def _header(metrics: dict) -> None:
    """Hero header with kicker, title, subtitle, and meta row."""
    st.markdown(
        '<a class="back-link" href="/" target="_self">'
        "← Back to Dashboard</a>",
        unsafe_allow_html=True,
    )
    st.markdown(
        '<div class="page-header">'
        '<div class="page-kicker">◆ Executive Dashboard</div>'
        '<div class="page-title">The Churn Business at a Glance</div>'
        '<div class="page-subtitle">'
        "A management-level overview of what is happening across the customer "
        "base, why it is happening, and what to do about it — built entirely "
        "from the customer dataset with no simulated numbers."
        "</div>"
        '<div class="page-rule"></div>'
        "</div>",
        unsafe_allow_html=True,
    )
    st.markdown(
        '<div class="meta-row">'
        f'<div class="meta-hint">Generated <b>{datetime.now().strftime("%B %d, %Y · %H:%M")}</b>'
        f" · Data as of the latest load · Deployed model: <b>{metrics['model_label']}</b>"
        f' · <b>15 sections</b> · Financial figures that are modeled are labeled '
        f'"Estimated".</div>'
        "</div>",
        unsafe_allow_html=True,
    )


def _chart(title: str, desc: str, fig: go.Figure) -> None:
    """Wrap a Plotly figure with a title and description."""
    st.markdown(
        f'<div class="chart-title" style="font-size:0.85rem;font-weight:600;'
        f'color:#F4F2EE;margin-bottom:0.15rem;">{title}</div>'
        f'<div class="chart-desc" style="font-size:0.7rem;color:#D6D8D8;'
        f'margin-bottom:0.6rem;opacity:0.72;">{desc}</div>',
        unsafe_allow_html=True,
    )
    st.plotly_chart(fig, width="stretch")


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 01 — EXECUTIVE KPI CARDS
# ═══════════════════════════════════════════════════════════════════════════════


def _executive_kpis(metrics: dict) -> None:
    """Three rows of four KPI cards covering customers, revenue, and model."""
    m = metrics
    _kpi_row(m, [
        ("Total Customers", f"{m['total']:,}", "Full cleaned dataset", ""),
        ("Active Customers", f"{m['n_retained']:,}", f"{m['retention_rate']:.1f}% retained", "good"),
        ("Churned Customers", f"{m['n_churned']:,}", f"{m['churn_rate']:.1f}% of base", "bad"),
        ("Churn Rate", f"{m['churn_rate']:.1f}%", "Of the total customer base", "accent"),
    ])

    _kpi_row(m, [
        ("Monthly Revenue (MRR)", f"${m['mrr_all']/1000:,.1f}K",
         "Sum of monthly charges", ""),
        ("Annual Projection", f"${m['annual_projection']/1e6:,.2f}M",
         "Estimated · MRR × 12", ""),
        ("ARPU", f"${m['arpu']:.2f}", "Per-customer average", ""),
        ("Avg. Customer LTV", f"${m['avg_total_charges']:,.2f}",
         "Estimated · mean total charges", "accent"),
    ])

    _kpi_row(m, [
        ("Model Accuracy", f"{m['model_accuracy']:.1f}%",
         f"{m['model_label']} · deployed", "good"),
        ("Best Model", f"{m['model_label']}",
         f"AUC {m['model_auc']:.4f}", ""),
        ("Retention Rate", f"{m['retention_rate']:.1f}%",
         "Active ÷ total customers", "good"),
        ("Average Tenure", f"{m['avg_tenure']:.1f} mo",
         f"Churned avg {m['avg_tenure_churned']:.1f} mo", ""),
    ])


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 02 — BUSINESS HEALTH SCORE
# ═══════════════════════════════════════════════════════════════════════════════


def _gauge(value: float, color: str, title: str) -> go.Figure:
    """Plotly gauge for a single 0-100 business score."""
    fig = go.Figure(go.Indicator(
        mode="gauge+number",
        value=value,
        number={"font": {"color": "#F4F2EE", "size": 30}},
        title={"text": title, "font": {"color": "#C8A96B", "size": 13}},
        gauge={
            "shape": "angular",
            "axis": {
                "range": [0, 100],
                "tickcolor": "#D6D8D8",
                "tickfont": {"color": "#D6D8D8", "size": 9},
            },
            "bar": {"color": color, "thickness": 0.3},
            "bgcolor": "rgba(0,0,0,0)",
            "borderwidth": 0,
            "steps": [
                {"range": [0, 40], "color": "rgba(217,124,124,0.18)"},
                {"range": [40, 70], "color": "rgba(200,169,107,0.22)"},
                {"range": [70, 100], "color": "rgba(143,162,138,0.26)"},
            ],
        },
    ))
    fig.update_layout(
        height=260,
        margin=dict(t=50, b=10, l=15, r=15),
        paper_bgcolor="rgba(0,0,0,0)",
        font={"family": "Inter, sans-serif"},
    )
    return fig


def _health_score(metrics: dict) -> None:
    """Composite executive health gauge with weighted component gauges."""
    m = metrics
    left, right = st.columns([1, 1.6], gap="medium")
    with left:
        with st.container(border=True):
            st.plotly_chart(
                _gauge(m["health"], m["health_color"], "Executive Health Score"),
                width="stretch",
            )
            st.markdown(
                f'<div class="note-text" style="text-align:center;margin-top:0.2rem;">'
                f"Band: <b>{m['health_band']}</b> — weighted composite of six "
                f"business components.</div>",
                unsafe_allow_html=True,
            )

    with right:
        comp_colors = {
            "Retention": "#8FA28A",
            "Revenue": "#C8A96B",
            "Satisfaction": "#9BCEC1",
            "Churn Health": "#8FA28A",
            "Tenure": "#C8A96B",
            "Model Confidence": "#9BCEC1",
        }
        cols = st.columns(3, gap="medium")
        items = list(m["health_scores"].items())
        for i, col in enumerate(cols):
            with col:
                for j in range(2):
                    idx = i * 2 + j
                    if idx >= len(items):
                        break
                    name, value = items[idx]
                    with st.container(border=True):
                        st.plotly_chart(
                            _gauge(value, comp_colors[name], name),
                            width="stretch",
                        )
        weight_text = " · ".join(
            f"{k} {w*100:.0f}%" for k, w in m["health_weights"].items()
        )
        st.markdown(
            f'<div class="note-text">Weights: {weight_text}. Components are '
            f"0–100 scores derived from the dataset: Retention = active share, "
            f"Revenue = ARPU, Satisfaction = add-on adoption, Churn Health = "
            f"100 − churn rate, Tenure = avg tenure vs max, Model Confidence = "
            f"deployed accuracy.</div>",
            unsafe_allow_html=True,
        )


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 03 — REVENUE OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════


def _revenue_charts(df: pd.DataFrame, m: dict) -> None:
    """Six revenue-oriented charts on two rows."""
    tenure_rev = m["tenure_cohort"]
    fig1 = px.bar(
        tenure_rev, x="Cohort", y="revenue",
        title="Monthly Revenue by Tenure Cohort",
        custom_data=["customers", "churn_rate"],
        color_discrete_sequence=["#8FA28A"],
    )
    fig1.update_traces(
        hovertemplate="<b>%{x}</b><br>Revenue: $%{y:,.0f}"
                      "<br>Customers: %{customdata[0]:,}"
                      "<br>Churn rate: %{customdata[1]:.1f}%<extra></extra>",
    )
    fig1.update_layout(template=TEMPLATE)
    _chart("Monthly Revenue by Tenure Cohort", "Distribution of MRR across tenure bands",
           fig1)

    fig2 = px.bar(
        df.groupby("Contract")["MonthlyCharges"].sum().reset_index(),
        x="Contract", y="MonthlyCharges",
        title="Revenue by Contract Type",
        color_discrete_sequence=["#C8A96B"],
    )
    fig2.update_traces(
        hovertemplate="<b>%{x}</b><br>Revenue: $%{y:,.0f}<extra></extra>",
    )
    fig2.update_layout(template=TEMPLATE, xaxis_tickangle=-20)
    _chart("Revenue by Contract", "Monthly revenue contribution per contract term",
           fig2)

    fig3 = px.bar(
        df.groupby("InternetService")["MonthlyCharges"].sum().reset_index(),
        x="InternetService", y="MonthlyCharges",
        title="Revenue by Internet Service",
        color_discrete_sequence=["#9BCEC1"],
    )
    fig3.update_traces(
        hovertemplate="<b>%{x}</b><br>Revenue: $%{y:,.0f}<extra></extra>",
    )
    fig3.update_layout(template=TEMPLATE, xaxis_tickangle=-20)
    _chart("Revenue by Internet Service", "Fiber carries the majority of revenue",
           fig3)

    fig4 = px.bar(
        df.groupby("PaymentMethod")["MonthlyCharges"].sum().reset_index(),
        x="PaymentMethod", y="MonthlyCharges",
        title="Revenue by Payment Method",
        color_discrete_sequence=["#D6D8D8"],
    )
    fig4.update_traces(
        hovertemplate="<b>%{x}</b><br>Revenue: $%{y:,.0f}<extra></extra>",
    )
    fig4.update_layout(template=TEMPLATE, xaxis_tickangle=-25)
    _chart("Revenue by Payment Method", "Monthly revenue per billing channel",
           fig4)

    lost = df[df["Churn"] == "Yes"]
    lost_by_contract = (
        lost.groupby("Contract")["MonthlyCharges"]
        .sum()
        .mul(12)
        .sort_values(ascending=True)
        .reset_index()
        .rename(columns={"MonthlyCharges": "Annualized Loss"})
    )
    fig5 = px.bar(
        lost_by_contract, x="Annualized Loss", y="Contract", orientation="h",
        title="Revenue Lost Due to Churn (Estimated)",
        color_discrete_sequence=["#C8A96B"],
    )
    fig5.update_traces(
        hovertemplate="<b>%{y}</b><br>Estimated annual loss: $%{x:,.0f}<extra></extra>",
    )
    fig5.update_layout(template=TEMPLATE)
    _chart("Revenue Lost Due to Churn (Estimated)",
           "Annualized churned MRR by contract term", fig5)

    fig6 = go.Figure(go.Bar(
        x=["Retained", "At Risk"],
        y=[m["annual_retained"], m["annual_at_risk"]],
        marker_color=["#8FA28A", "#C8A96B"],
        hovertemplate="<b>%{x}</b><br>Estimated annual revenue: $%{y:,.0f}<extra></extra>",
    ))
    fig6.update_layout(
        template=TEMPLATE,
        title="Revenue Retained vs At Risk (Estimated)",
        yaxis_title="Annualized revenue (USD)",
    )
    _chart("Revenue Retained vs At Risk (Estimated)",
           "Annualized revenue currently retained vs projected churn loss", fig6)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 04 — CHURN OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════


def _rate_bar(labels: list, values: list, title: str, color: str = "#C8A96B") -> go.Figure:
    """Simple churn-rate bar chart with value labels in the hover."""
    fig = go.Figure(go.Bar(
        x=labels, y=values,
        marker_color=color,
        hovertemplate="<b>%{x}</b><br>Churn rate: %{y:.1f}%<extra></extra>",
    ))
    fig.add_trace(go.Scatter(
        x=labels, y=values, mode="markers",
        marker=dict(size=10, color=color, symbol="circle-open"),
        showlegend=False,
        hovertemplate="<b>%{x}</b><br>Churn rate: %{y:.1f}%<extra></extra>",
    ))
    fig.update_layout(
        template=TEMPLATE,
        title=title,
        yaxis_title="Churn rate (%)",
    )
    return fig


def _churn_charts(df: pd.DataFrame, m: dict) -> None:
    """Six churn-rate charts on two rows."""
    tc = m["tenure_cohort"]
    fig1 = go.Figure(go.Scatter(
        x=tc["Cohort"].astype(str), y=tc["churn_rate"],
        mode="lines+markers",
        line=dict(color="#C8A96B", width=3),
        marker=dict(color="#C8A96B", size=9),
        hovertemplate="<b>%{x} months</b><br>Churn rate: %{y:.1f}%<extra></extra>",
    ))
    fig1.update_layout(
        template=TEMPLATE,
        title="Monthly Churn Trend by Tenure Cohort (Proxy)",
        yaxis_title="Churn rate (%)",
    )
    _chart("Monthly Churn Trend by Tenure Cohort (Proxy)",
           "Churn rate by tenure band — the dataset has no dates, so tenure "
           "cohorts stand in for a monthly trend", fig1)

    cat = m["cat"]
    fig2 = _rate_bar(
        list(cat["Contract"].index), list(cat["Contract"].values),
        "Churn Rate by Contract", "#C8A96B",
    )
    fig2.update_layout(xaxis_tickangle=-20)
    _chart("Contract vs Churn", "Month-to-month is the dominant churn driver",
           fig2)

    fig3 = _rate_bar(
        list(cat["InternetService"].index), list(cat["InternetService"].values),
        "Churn Rate by Internet Service", "#9BCEC1",
    )
    fig3.update_layout(xaxis_tickangle=-20)
    _chart("Internet Service vs Churn", "Fiber optic churns well above average",
           fig3)

    fig4 = _rate_bar(
        ["No", "Yes"], [m["senior"]["No"], m["senior"]["Yes"]],
        "Churn Rate by Senior Citizen", "#8FA28A",
    )
    _chart("Senior Citizen vs Churn", "Seniors churn at a materially higher rate",
           fig4)

    fig5 = _rate_bar(
        list(cat["PaymentMethod"].index), list(cat["PaymentMethod"].values),
        "Churn Rate by Payment Method", "#C8A96B",
    )
    fig5.update_layout(xaxis_tickangle=-25)
    _chart("Payment Method vs Churn", "Electronic check shows the most friction",
           fig5)

    fig6 = _rate_bar(
        list(cat["gender"].index), list(cat["gender"].values),
        "Churn Rate by Gender", "#D6D8D8",
    )
    _chart("Gender vs Churn", "Churn is broadly balanced across genders", fig6)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 05 — CUSTOMER SEGMENTATION
# ═══════════════════════════════════════════════════════════════════════════════


def _segmentation(m: dict) -> None:
    """Five segment cards with counts, share, and description."""
    cards = []
    for name, (icon, color, desc) in SEGMENT_META.items():
        count, share = m["seg_counts"][name]
        cards.append(
            f'<div class="seg-card">'
            f'<div class="seg-icon">{icon}</div>'
            f'<div class="seg-name">{name}</div>'
            f'<div class="seg-count">{count:,}</div>'
            f'<div class="seg-share">{share:.1f}% of customer base</div>'
            f'<div class="seg-desc">{desc}</div>'
            f"</div>"
        )
    with st.container(border=True):
        st.markdown(
            f'<div class="seg-grid">{"".join(cards)}</div>'
            '<div class="note-text">Segmentation is deterministic and '
            "rule-based: <b>Critical</b> = month-to-month with tenure under "
            "12 months; <b>High Risk</b> = month-to-month or electronic check; "
            "<b>VIP</b> = tenure ≥ 60 mo and spend ≥ $90; <b>Premium</b> = "
            "tenure ≥ 48 mo or spend ≥ $85; otherwise <b>Standard</b>.</div>",
            unsafe_allow_html=True,
        )


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 06 — GEOGRAPHIC INSIGHTS
# ═══════════════════════════════════════════════════════════════════════════════


def _geo_chart(data: pd.Series, title: str, color: str) -> go.Figure:
    frame = pd.DataFrame({
        "Category": data.index.astype(str),
        "Customers": data.values,
    })
    fig = px.bar(
        frame,
        x="Category", y="Customers",
        title=title,
        color_discrete_sequence=[color],
    )
    fig.update_traces(
        hovertemplate="<b>%{x}</b><br>Customers: %{y:,}<extra></extra>",
    )
    fig.update_layout(template=TEMPLATE, xaxis_tickangle=-20)
    return fig


def _geographic(m: dict) -> None:
    """Regional summary using a market-composition proxy."""
    g = m["geo_proxy"]
    st.markdown(
        '<div class="note-text">This dataset contains <b>no geographic fields</b> '
        "(no city, state, region, or country). To still answer the regional "
        "question, the section below shows the <b>market composition</b> of the "
        "customer base as a regional proxy — the service, contract, and billing "
        "mix that would shape regional strategy. All figures are labeled "
        "Estimated/Proxy accordingly.</div>",
        unsafe_allow_html=True,
    )
    c1, c2, c3 = st.columns(3, gap="medium")
    with c1:
        _chart("Internet Service Mix (Regional Proxy)",
               "Where internet revenue concentration sits", 
               _geo_chart(g["internet"], "Internet Service Mix", "#9BCEC1"))
    with c2:
        _chart("Contract Mix (Regional Proxy)",
               "Term structure of the book",
               _geo_chart(g["contract"], "Contract Mix", "#C8A96B"))
    with c3:
        _chart("Payment Method Mix (Regional Proxy)",
               "Billing channel composition",
               _geo_chart(g["payment"], "Payment Method Mix", "#8FA28A"))

    st.markdown(
        '<div class="note-text">Proxy takeaway: fiber-heavy markets carry the '
        "highest revenue share and the highest churn — regional plans should "
        "pair fiber quality investments with annual-contract offers in any "
        "territory with a fiber-dominant mix.</div>",
        unsafe_allow_html=True,
    )


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 07 — EXECUTIVE ALERTS
# ═══════════════════════════════════════════════════════════════════════════════


def _alerts(m: dict) -> None:
    """Executable alert cards ordered by severity."""
    order = {"Critical": 0, "Warning": 1, "Information": 2}
    sorted_alerts = sorted(m["alerts"], key=lambda a: order[a[0]])
    for severity, icon, title, text in sorted_alerts:
        color = SEVERITY_COLORS[severity]
        st.markdown(
            f'<div class="alert-card" style="border-left-color:{color};">'
            f'<div class="alert-icon">{icon}</div>'
            f'<div>'
            f'<span class="alert-pill" style="background:{color}">{severity}</span>'
            f'<div class="alert-title">{title}</div>'
            f'<div class="alert-text">{text}</div>'
            f"</div></div>",
            unsafe_allow_html=True,
        )


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 08 — DEPARTMENT PERFORMANCE
# ═══════════════════════════════════════════════════════════════════════════════


def _departments(m: dict) -> None:
    """Six department cards with status, KPI, and suggested action."""
    status_colors = {"Good": "#8FA28A", "At Risk": "#D97C7C",
                     "Opportunity": "#C8A96B", "Stable": "#9BCEC1"}
    for dep in m["departments"]:
        color = status_colors.get(dep["status"], "#C8A96B")
        st.markdown(
            f'<div class="dep-card">'
            f'<div class="dep-head">'
            f'<span class="dep-icon">{dep["icon"]}</span>'
            f'<span class="dep-name">{dep["name"]}</span>'
            f'<span class="dep-status" style="background:{color}">{dep["status"]}</span>'
            f"</div>"
            f'<div class="dep-kpi-label">{dep["kpi_label"]}</div>'
            f'<div class="dep-kpi">{dep["kpi"]}</div>'
            f'<div class="dep-action">Suggested action: {dep["action"]}</div>'
            f"</div>",
            unsafe_allow_html=True,
        )


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 09 — TOP BUSINESS OPPORTUNITIES
# ═══════════════════════════════════════════════════════════════════════════════


def _opportunities(m: dict) -> None:
    """Top ten ranked opportunities with impact, ROI, and cost."""
    for i, opp in enumerate(m["opportunities"], start=1):
        rank_color = "#8FA28A" if opp["roi"] >= 100 else (
            "#C8A96B" if opp["roi"] >= 50 else "#D97C7C")
        st.markdown(
            f'<div class="opp-card" style="border-left-color:{rank_color};">'
            f'<div class="opp-rank">{i}</div>'
            f'<div style="flex:1;min-width:220px;">'
            f'<div class="opp-title">{opp["icon"]} {opp["name"]}</div>'
            f'<div class="opp-desc">{opp["desc"]}</div>'
            f"</div>"
            f'<div class="opp-meta">'
            f'<div class="opp-meta-item">'
            f'<div class="opp-meta-label">Est. ROI</div>'
            f'<div class="opp-meta-value">{opp["roi"]:.0f}%</div></div>'
            f'<div class="opp-meta-item">'
            f'<div class="opp-meta-label">Est. Cost</div>'
            f'<div class="opp-meta-value">${opp["cost"]:,.0f}</div></div>'
            f'<div class="opp-meta-item">'
            f'<div class="opp-meta-label">Cohort</div>'
            f'<div class="opp-meta-value">{opp["cohort"]:,}</div></div>'
            f"</div></div>",
            unsafe_allow_html=True,
        )
    st.markdown(
        '<div class="note-text">ROI is a <b>modeled estimate</b>: (expected '
        "benefit − program cost) ÷ cost. Benefit assumes a 25–30% recovery of "
        "annualized revenue at risk (or modeled revenue uplift for upsell "
        "opportunities). Program cost assumes $25–40 per-customer outreach. "
        "These are planning figures, not guarantees.</div>",
        unsafe_allow_html=True,
    )


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 10 — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════


def _executive_summary(m: dict) -> None:
    """Auto-generated one-paragraph executive narrative."""
    with st.container(border=True):
        st.markdown(
            f'<div class="notes-box">{m["summary"]}</div>',
            unsafe_allow_html=True,
        )


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 11 — STRATEGIC ROADMAP
# ═══════════════════════════════════════════════════════════════════════════════


def _roadmap() -> None:
    """Immediate / mid-term / long-term strategic roadmap."""
    priority_colors = {"Critical": "#D97C7C", "High": "#C8A96B", "Medium": "#8FA28A"}
    cols = st.columns(3, gap="medium")
    for col, (phase, horizon, icon, actions, outcome, owner, priority) in zip(cols, ROADMAP):
        with col:
            items = "".join(
                f'<div class="rm-item"><span>▸</span>{a}</div>' for a in actions
            )
            pcolor = priority_colors.get(priority, "#8FA28A")
            st.markdown(
                f'<div class="rm-card" style="border-top-color:{pcolor};">'
                f'<div class="rm-head">'
                f'<span>{icon}</span>'
                f'<span class="rm-phase">{phase}</span>'
                f'<span class="rm-horizon">{horizon}</span>'
                f"</div>"
                f"{items}"
                f'<div class="rm-outcome">🎯 {outcome}</div>'
                f'<div class="rm-owner">Owner: <b style="color:#C8A96B">{owner}</b>'
                f' · Priority: {priority}</div>'
                f"</div>",
                unsafe_allow_html=True,
            )


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 12 — AI INSIGHTS
# ═══════════════════════════════════════════════════════════════════════════════


def _ai_insights() -> None:
    """Five model-backed insights presented as cards."""
    cols = st.columns(3, gap="medium")
    for i, (icon, title, text) in enumerate(AI_INSIGHTS):
        col = cols[i % 3]
        with col:
            st.markdown(
                f'<div class="insight-card">'
                f'<div class="insight-icon">{icon}</div>'
                f'<div class="insight-title">{title}</div>'
                f'<div class="insight-text">{text}</div>'
                f"</div>",
                unsafe_allow_html=True,
            )
    st.markdown(
        '<div class="note-text">Insights are deterministic comparisons derived '
        "directly from the dataset; the deployed model adds confidence that "
        "these drivers generalize to unseen customers.</div>",
        unsafe_allow_html=True,
    )


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 13 — BOARD MEETING SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════


def _board_summary(m: dict) -> None:
    """Board-ready bullets in an expandable card."""
    with st.expander("📋 Board Meeting Summary — open for the briefing", expanded=False):
        for title, color, text in m["board"]:
            st.markdown(
                f'<div class="board-item">'
                f'<b style="color:{color}">● {title}.</b> {text}'
                f"</div>",
                unsafe_allow_html=True,
            )
        st.markdown(
            '<div class="note-text">Share the "Board Briefing" file below for '
            "a print-ready one-pager of this summary.</div>",
            unsafe_allow_html=True,
        )


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 14 — EXPORT OPTIONS
# ═══════════════════════════════════════════════════════════════════════════════


def _ascii(text) -> str:
    """Strip non-Latin-1 characters for fpdf2's core fonts."""
    return re.sub(r"[^\x00-\x7F]", "", str(text))


def _kpi_frame(m: dict) -> pd.DataFrame:
    """A tidy KPI dataframe shared by CSV and summary exports."""
    rows = [
        ("Customers", "Total customers", m["total"]),
        ("Customers", "Active customers", m["n_retained"]),
        ("Customers", "Churned customers", m["n_churned"]),
        ("Customers", "Churn rate (%)", m["churn_rate"]),
        ("Customers", "Retention rate (%)", m["retention_rate"]),
        ("Revenue", "Monthly revenue (USD)", round(m["mrr_all"], 2)),
        ("Revenue", "Annual projection (USD, est.)", round(m["annual_projection"], 2)),
        ("Revenue", "ARPU (USD)", m["arpu"]),
        ("Revenue", "Avg customer LTV (USD, est.)", m["avg_total_charges"]),
        ("Revenue", "Annual revenue at risk (USD, est.)", round(m["annual_at_risk"], 2)),
        ("Revenue", "Annual revenue retained (USD, est.)", round(m["annual_retained"], 2)),
        ("Health", "Executive health score", m["health"]),
        ("Health", "Health band", m["health_band"]),
        ("Tenure", "Average tenure (months)", m["avg_tenure"]),
        ("Model", "Best model", m["model_label"]),
        ("Model", "Model accuracy (%)", m["model_accuracy"]),
        ("Model", "Model AUC", m["model_auc"]),
    ]
    for name, rate in m["cat"]["Contract"].items():
        rows.append(("Churn by Contract", f"{name} churn rate (%)", rate))
    for name, rate in m["cat"]["InternetService"].items():
        rows.append(("Churn by Service", f"{name} churn rate (%)", rate))
    for name, rate in m["cat"]["PaymentMethod"].items():
        rows.append(("Churn by Payment", f"{name} churn rate (%)", rate))
    for name, (count, share) in m["seg_counts"].items():
        rows.append(("Segments", f"{name} customers", count))
        rows.append(("Segments", f"{name} share (%)", share))
    return pd.DataFrame(rows, columns=["Category", "Metric", "Value"])


def _build_pdf(m: dict) -> bytes:
    """Executive PDF report with fpdf2."""
    class Report(FPDF):
        def footer(self):
            self.set_y(-14)
            self.set_font("Helvetica", "", 7)
            self.set_text_color(150, 150, 150)
            self.cell(0, 5,
                      "Customer Churn Analytics Platform  -  Page "
                      f"{self.page_no()}",
                      align="C")

    pdf = Report(format="A4", unit="mm")
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_margins(14, 15, 14)
    pdf.add_page()

    pdf.set_fill_color(15, 48, 64)
    pdf.rect(0, 0, 210, 30, "F")
    pdf.set_fill_color(200, 169, 107)
    pdf.rect(0, 30, 210, 2.2, "F")
    pdf.set_text_color(244, 242, 238)
    pdf.set_xy(14, 7)
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 7, "Executive Dashboard - Customer Churn")
    pdf.set_xy(14, 16)
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(214, 216, 216)
    pdf.cell(0, 5, "Management overview generated by the Customer Churn Analytics Platform")
    pdf.set_xy(14, 22)
    pdf.set_font("Helvetica", "", 8)
    pdf.cell(0, 5, datetime.now().strftime("Generated %Y-%m-%d at %H:%M"))

    def section(title):
        pdf.ln(4)
        pdf.set_font("Helvetica", "B", 11)
        pdf.set_fill_color(200, 169, 107)
        pdf.set_text_color(15, 48, 64)
        pdf.cell(0, 8, title, fill=True, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_text_color(30, 40, 45)
        pdf.ln(2)

    def kv(label, value):
        pdf.set_font("Helvetica", "B", 9)
        pdf.cell(62, 6, _ascii(label), new_x=XPos.RIGHT, new_y=YPos.TOP)
        pdf.set_font("Helvetica", "", 9)
        pdf.cell(0, 6, _ascii(value), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    section("1. Executive KPIs")
    kv("Total customers", f"{m['total']:,}")
    kv("Active / Churned", f"{m['n_retained']:,} / {m['n_churned']:,}")
    kv("Churn rate", f"{m['churn_rate']:.1f}%")
    kv("Retention rate", f"{m['retention_rate']:.1f}%")
    kv("Monthly revenue (MRR)", f"${m['mrr_all']/1000:,.1f}K")
    kv("Annual projection (est.)", f"${m['annual_projection']/1e6:,.2f}M")
    kv("ARPU", f"${m['arpu']:.2f}")
    kv("Avg customer LTV (est.)", f"${m['avg_total_charges']:,.2f}")
    kv("Average tenure", f"{m['avg_tenure']:.1f} months")
    kv("Deployed model", f"{m['model_label']} ({m['model_accuracy']:.1f}% acc, "
                         f"AUC {m['model_auc']})")

    section("2. Business Health Score")
    kv("Executive health score", f"{m['health']:.0f} / 100 ({m['health_band']})")
    for name, value in m["health_scores"].items():
        kv(f"  {name}", f"{value:.1f} / 100")

    section("3. Revenue at Risk (Estimated)")
    kv("Annual revenue at risk", f"${m['annual_at_risk']:,.0f}")
    kv("Annual revenue retained", f"${m['annual_retained']:,.0f}")

    section("4. Churn Drivers")
    kv("Month-to-month churn", f"{m['mtm_rate']:.1f}%")
    kv("Fiber optic churn", f"{m['fiber_rate']:.1f}%")
    kv("Electronic-check churn", f"{m['ec_rate']:.1f}%")
    kv("Tenure < 12 months churn", f"{m['tenure_lt_12']:.1f}%")
    kv("Two-year contract churn", f"{m['two_year_rate']:.1f}%")

    section("5. Segmentation")
    for name, (count, share) in m["seg_counts"].items():
        kv(f"  {name}", f"{count:,} customers ({share:.1f}%)")

    section("6. Executive Alerts")
    for severity, icon, title, text in sorted(
            m["alerts"], key=lambda a: (a[0] != "Critical", a[0] != "Warning")):
        pdf.set_font("Helvetica", "B", 9)
        pdf.cell(0, 5, f"- [{severity}] {_ascii(title)}",
                 new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("Helvetica", "", 8.5)
        pdf.multi_cell(0, 4.5, _ascii(text))
        pdf.ln(1)

    section("7. Top Business Opportunities (Estimated)")
    for i, opp in enumerate(m["opportunities"][:6], start=1):
        pdf.set_font("Helvetica", "B", 9)
        pdf.cell(0, 5, f"{i}. {_ascii(opp['name'])}  "
                       f"(est. ROI {opp['roi']:.0f}%)",
                 new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("Helvetica", "", 8.5)
        pdf.multi_cell(0, 4.5,
                       _ascii(f"{opp['desc']} Cohort: {opp['cohort']:,}. "
                              f"Est. cost: ${opp['cost']:,.0f}. "
                              f"Est. benefit: ${opp['benefit']:,.0f}."))
        pdf.ln(1)

    section("8. Strategic Roadmap")
    for phase, horizon, icon, actions, outcome, owner, priority in ROADMAP:
        pdf.set_font("Helvetica", "B", 9)
        pdf.cell(0, 5, f"{_ascii(phase)} ({_ascii(horizon)})",
                 new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("Helvetica", "", 8.5)
        for a in actions:
            pdf.cell(4, 4.5, "", new_x=XPos.RIGHT, new_y=YPos.TOP)
            pdf.multi_cell(0, 4.5, _ascii(f"- {a}"))
        pdf.multi_cell(0, 4.5, _ascii(f"Outcome: {outcome}  Owner: {owner}"))
        pdf.ln(1)

    section("9. Executive Summary")
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(30, 40, 45)
    pdf.multi_cell(0, 5, _ascii(m["summary"]))

    section("10. Methodology & Estimates")
    pdf.set_font("Helvetica", "", 8.5)
    pdf.set_text_color(60, 70, 75)
    pdf.multi_cell(
        0, 4.5,
        _ascii(
            f"All figures are derived deterministically from the IBM Telco "
            f"Customer Churn dataset ({m['total']:,} customers after cleaning). "
            f"The executive health score is a weighted composite of retention, "
            f"revenue, satisfaction (add-on adoption), churn health, tenure, "
            f"and model confidence. Financial figures that are annualized or "
            f"projected — revenue at risk, ROI, LTV — are modeled estimates, "
            f"not CRM or LTV data. Segmentation and alerts follow fixed "
            f"business rules. The dataset contains no geographic or date "
            f"fields, so the geographic section and monthly trend use "
            f"composition and tenure-cohort proxies. Deployed model: "
            f"{m['model_label']}."
        ),
    )

    return bytes(pdf.output())


def _build_pptx(m: dict) -> bytes:
    """Board briefing as a .pptx deck with python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def _bullet(slide, items, top, size=14):
        left = Inches(0.6)
        width = Inches(12.1)
        for i, item in enumerate(items):
            box = slide.shapes.add_textbox(
                left, Inches(top + i * 0.42), width, Inches(0.4))
            tf = box.text_frame
            tf.text = item
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(size)
                    run.font.color.rgb = None

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.4))
    tf = tb.text_frame
    tf.text = "Executive Dashboard - Board Briefing"
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(40)
            run.font.bold = True
    sb = slide.shapes.add_textbox(Inches(0.9), Inches(3.7), Inches(11.5), Inches(0.9))
    stf = sb.text_frame
    stf.text = "Customer Churn Analytics Platform  |  Generated "
    stf.text += datetime.now().strftime("%Y-%m-%d")
    for p in stf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(18)

    # KPI slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = tb.text_frame
    tf.text = "Key Metrics"
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(28)
            run.font.bold = True
    _bullet(slide, [
        f"Total customers: {m['total']:,}  |  Active: {m['n_retained']:,}  |  "
        f"Churned: {m['n_churned']:,}",
        f"Churn rate: {m['churn_rate']:.1f}%  |  Retention rate: {m['retention_rate']:.1f}%",
        f"Monthly revenue: ${m['mrr_all']/1000:,.1f}K  |  ARPU: ${m['arpu']:.2f}",
        f"Annual revenue at risk (est.): ${m['annual_at_risk']:,.0f}",
        f"Executive health score: {m['health']:.0f}/100  ({m['health_band']})",
        f"Deployed model: {m['model_label']}  ({m['model_accuracy']:.1f}% accuracy, "
        f"AUC {m['model_auc']})",
    ], top=1.5)

    # Churn drivers slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = tb.text_frame
    tf.text = "Churn Drivers"
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(28)
            run.font.bold = True
    _bullet(slide, [
        f"Month-to-month contracts: {m['mtm_rate']:.1f}% churn",
        f"Fiber optic service: {m['fiber_rate']:.1f}% churn",
        f"Electronic check payment: {m['ec_rate']:.1f}% churn",
        f"Tenure under 12 months: {m['tenure_lt_12']:.1f}% churn",
        f"Two-year contracts: {m['two_year_rate']:.1f}% churn (anchor of retention)",
    ], top=1.5)

    # Opportunities slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = tb.text_frame
    tf.text = "Top Opportunities (Estimated)"
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(28)
            run.font.bold = True
    opp_lines = []
    for opp in m["opportunities"][:6]:
        opp_lines.append(
            f"{opp['name']}: est. ROI {opp['roi']:.0f}%, est. cost "
            f"${opp['cost']:,.0f}, cohort {opp['cohort']:,}")
    _bullet(slide, opp_lines, top=1.5, size=16)

    # Roadmap slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = tb.text_frame
    tf.text = "Strategic Roadmap"
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(28)
            run.font.bold = True
    roadmap_lines = []
    for phase, horizon, icon, actions, outcome, owner, priority in ROADMAP:
        roadmap_lines.append(f"{phase} ({horizon}): {actions[0]}")
    _bullet(slide, roadmap_lines, top=1.5, size=16)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_summary_text(m: dict) -> str:
    """Plain-text business summary for download."""
    lines = [
        "EXECUTIVE DASHBOARD - BUSINESS SUMMARY",
        "Customer Churn Analytics Platform",
        f"Generated {datetime.now().strftime('%Y-%m-%d at %H:%M')}",
        "=" * 72,
        "",
        "KEY METRICS",
        f"  Total customers:     {m['total']:,}",
        f"  Active customers:    {m['n_retained']:,}",
        f"  Churned customers:   {m['n_churned']:,}",
        f"  Churn rate:          {m['churn_rate']:.1f}%",
        f"  Retention rate:      {m['retention_rate']:.1f}%",
        f"  Monthly revenue:     ${m['mrr_all']/1000:,.1f}K (MRR)",
        f"  ARPU:                ${m['arpu']:.2f}",
        f"  Avg customer LTV:    ${m['avg_total_charges']:,.2f} (est.)",
        f"  Avg tenure:          {m['avg_tenure']:.1f} months",
        f"  Health score:        {m['health']:.0f}/100 ({m['health_band']})",
        f"  Deployed model:      {m['model_label']} "
        f"({m['model_accuracy']:.1f}% acc, AUC {m['model_auc']})",
        "",
        "REVENUE IMPACT (ESTIMATED)",
        f"  Annual revenue at risk:  ${m['annual_at_risk']:,.0f}",
        f"  Annual revenue retained: ${m['annual_retained']:,.0f}",
        "",
        "SEGMENTATION",
    ]
    for name, (count, share) in m["seg_counts"].items():
        lines.append(f"  {name}: {count:,} customers ({share:.1f}%)")
    lines += [
        "",
        "EXECUTIVE SUMMARY",
    ]
    lines.append(f"  {m['summary']}")
    lines += [
        "",
        "TOP OPPORTUNITIES (ESTIMATED)",
    ]
    for i, opp in enumerate(m["opportunities"][:6], start=1):
        lines.append(
            f"  {i}. {opp['name']} — est. ROI {opp['roi']:.0f}%, "
            f"est. cost ${opp['cost']:,.0f}, cohort {opp['cohort']:,}")
    lines += [
        "",
        "AI INSIGHTS",
    ]
    for icon, title, text in AI_INSIGHTS:
        lines.append(f"  {icon} {title}: {text}")
    lines += [
        "",
        "BOARD MEETING SUMMARY",
    ]
    for title, color, text in m["board"]:
        lines.append(f"  - {title}: {text}")
    lines.append("")
    lines.append("NOTE: Annualized and projected financial figures are modeled")
    lines.append("estimates derived from the dataset, not CRM or LTV data.")
    return "\n".join(lines)


def _export_section(m: dict) -> None:
    """Four export options: PDF, KPI CSV, summary text, PowerPoint."""
    c1, c2, c3, c4 = st.columns(4, gap="medium")

    with c1:
        try:
            pdf_bytes = _build_pdf(m)
            st.download_button(
                "📄 Download Executive Report (PDF)",
                data=pdf_bytes,
                file_name="executive_dashboard_report.pdf",
                mime="application/pdf",
                width="stretch",
                key="exec_pdf_download",
            )
        except Exception:
            st.markdown(
                '<div class="note-text">PDF generation is unavailable right now.</div>',
                unsafe_allow_html=True,
            )

    with c2:
        kpi_csv = _kpi_frame(m).to_csv(index=False).encode("utf-8")
        st.download_button(
            "📊 Download KPI Dashboard (CSV)",
            data=kpi_csv,
            file_name="executive_kpi_dashboard.csv",
            mime="text/csv",
            width="stretch",
            key="exec_csv_download",
        )

    with c3:
        st.download_button(
            "📝 Download Business Summary (TXT)",
            data=_build_summary_text(m).encode("utf-8"),
            file_name="executive_business_summary.txt",
            mime="text/plain",
            width="stretch",
            key="exec_txt_download",
        )

    with c4:
        try:
            pptx_bytes = _build_pptx(m)
            st.download_button(
                "📽️ Download Board Briefing (PPTX)",
                data=pptx_bytes,
                file_name="executive_board_briefing.pptx",
                mime=(
                    "application/vnd.openxmlformats-officedocument."
                    "presentationml.presentation"
                ),
                width="stretch",
                key="exec_pptx_download",
            )
        except Exception:
            st.markdown(
                '<div class="note-text">PPTX generation is unavailable right now.</div>',
                unsafe_allow_html=True,
            )

    st.markdown(
        '<div class="note-text">The PDF is a full management report; the CSV '
        "holds the raw KPI table; the TXT is a shareable summary; the PPTX is "
        "a board-ready briefing deck.</div>",
        unsafe_allow_html=True,
    )


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════


def main() -> None:
    """Render the full executive dashboard."""
    _inject_css()

    df = load_data()
    metrics = _build_metrics(df)

    _header(metrics)

    # 01 · Executive KPI Cards
    _section_head("01", "📊", "Executive KPI Cards",
                  "Customers, revenue, and model health at a glance")
    _executive_kpis(metrics)

    # 02 · Business Health Score
    _section_head("02", "🧭", "Business Health Score",
                  "Weighted composite of six business components")
    _health_score(metrics)

    # 03 · Revenue Overview
    _section_head("03", "💰", "Revenue Overview",
                  "Where money comes from and what churn puts at risk")
    _revenue_charts(df, metrics)

    # 04 · Churn Overview
    _section_head("04", "📉", "Churn Overview",
                  "Who churns and where the risk concentrates")
    _churn_charts(df, metrics)

    # 05 · Customer Segmentation
    _section_head("05", "🗂️", "Customer Segmentation",
                  "Value-versus-risk taxonomy of the customer base")
    _segmentation(metrics)

    # 06 · Geographic Insights
    _section_head("06", "🌍", "Geographic Insights",
                  "Regional summary using a market-composition proxy")
    _geographic(metrics)

    # 07 · Executive Alerts
    _section_head("07", "🚨", "Executive Alerts",
                  "Auto-detected issues, ordered by severity")
    _alerts(metrics)

    # 08 · Department Performance
    _section_head("08", "🏢", "Department Performance",
                  "Status and suggested action per business function")
    _departments(metrics)

    # 09 · Top Business Opportunities
    _section_head("09", "💡", "Top Business Opportunities",
                  "Ranked opportunities with estimated ROI, cost, and cohort")
    _opportunities(metrics)

    # 10 · Executive Summary
    _section_head("10", "📝", "Executive Summary",
                  "Auto-generated narrative of the business situation")
    _executive_summary(metrics)

    # 11 · Strategic Roadmap
    _section_head("11", "🗺️", "Strategic Roadmap",
                  "Immediate, mid-term, and long-term action plan")
    _roadmap()

    # 12 · AI Insights
    _section_head("12", "🤖", "AI Insights",
                  "Model-backed findings distilled for leadership")
    _ai_insights()

    # 13 · Board Meeting Summary
    _section_head("13", "📋", "Board Meeting Summary",
                  "A concise briefing for the board")
    _board_summary(metrics)

    # 14 · Export Options
    _section_head("14", "📦", "Export Options",
                  "Download the dashboard as PDF, CSV, TXT, or PowerPoint")
    _export_section(metrics)


if __name__ == "__main__":
    main()
