"""
pptx_report.py — Professional PowerPoint (.pptx) export.

Builds a presentation-ready executive deck in the dashboard's corporate
visual language (navy / gold) entirely in memory:

  1. Title
  2. Executive Summary
  3. Customer Details
  4. KPI Summary
  5. Business Recommendations
  6. Conclusion & Next Steps

Presentation layer only — reads the prediction result, the customer inputs,
and the recommendation-engine analysis; it never recomputes or alters any
prediction, SHAP, or business metric. Nothing is written to disk; the deck
is returned as bytes for st.download_button().
"""

from __future__ import annotations

import os
import sys
from datetime import datetime
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

try:
    import prediction
except ModuleNotFoundError:
    _dashboard_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    if _dashboard_dir not in sys.path:
        sys.path.insert(0, _dashboard_dir)
    import prediction

try:
    import report
except ModuleNotFoundError:
    _dashboard_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    if _dashboard_dir not in sys.path:
        sys.path.insert(0, _dashboard_dir)
    import report

# ── Brand palette (matches dashboard/theme.py / theme.css) ──────────────────────

PALETTE = {
    "navy": "0F3040",
    "navy_deep": "0A2632",
    "surface": "163949",
    "card": "234556",
    "gold": "C8A96B",
    "gold_bright": "D4B678",
    "gold_dark": "B09055",
    "sage": "8FA28A",
    "teal": "9BCEC1",
    "red": "D97C7C",
    "red_bright": "E0635A",
    "blue": "6EA8FE",
    "green": "5FCE8B",
    "emerald": "3FD6C0",
    "purple": "A78BFA",
    "white": "FFFFFF",
    "ink": "F4F2EE",
    "sub": "D6D8D8",
    "text": "1E282D",
    "body": "424C52",
    "muted": "788288",
    "light": "F3F4F4",
    "line": "DDE0E2",
}

RISK_COLOR = {
    "Low Risk": "sage",
    "Medium Risk": "gold_dark",
    "High Risk": "red",
}

VERDICT_COLOR = {
    "Likely to Stay": "sage",
    "Likely to Churn": "red",
}

FONT = "Calibri"

SLIDE_W = 13.333
SLIDE_H = 7.5


def _rgb(name: str) -> RGBColor:
    return RGBColor.from_string(PALETTE[name])


def _safe(value) -> str:
    return str(value) if value is not None else ""


def _money(value) -> str:
    try:
        return f"${float(value):,.0f}"
    except (TypeError, ValueError):
        return "$0"


# ── Low-level drawing helpers ───────────────────────────────────────────────────


def _rect(slide, x, y, w, h, fill=None, line=None, rounded=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
    return shape


def _tblock(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP,
            align=PP_ALIGN.LEFT, font=FONT):
    """A text box with one or more paragraphs.

    `paras` is a list of dicts: text, size, bold, color, italic,
    align, space_after, space_before, line_spacing.
    """
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, pa in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pa.get("align", align)
        if pa.get("space_after") is not None:
            p.space_after = Pt(pa["space_after"])
        if pa.get("space_before") is not None:
            p.space_before = Pt(pa["space_before"])
        if pa.get("line_spacing") is not None:
            p.line_spacing = pa["line_spacing"]
        run = p.add_run()
        run.text = _safe(pa.get("text", ""))
        run.font.name = pa.get("font", font)
        run.font.size = Pt(pa.get("size", 14))
        run.font.bold = pa.get("bold", False)
        run.font.italic = pa.get("italic", False)
        run.font.color.rgb = _rgb(pa.get("color", "text"))
    return box


def _card_text(shape, text, size, bold, color, align=PP_ALIGN.CENTER,
               anchor=MSO_ANCHOR.MIDDLE, font=FONT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _safe(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _header(slide, title, subtitle=""):
    _rect(slide, 0, 0, SLIDE_W, 1.0, fill="navy")
    _rect(slide, 0, 1.0, SLIDE_W, 0.06, fill="gold")
    _tblock(slide, 0.55, 0.2, 12.2, 0.6,
            [{"text": title, "size": 25, "bold": True, "color": "ink"}])
    if subtitle:
        _tblock(slide, 0.55, 0.62, 12.2, 0.3,
                [{"text": subtitle, "size": 11.5, "color": "sub"}])


def _footer(slide, index: int, total: int):
    _rect(slide, 0, SLIDE_H - 0.42, SLIDE_W, 0.42, fill="light")
    _rect(slide, 0, SLIDE_H - 0.42, SLIDE_W, 0.02, fill="gold")
    _tblock(slide, 0.55, SLIDE_H - 0.38, 8.5, 0.3,
            [{"text": "Customer Churn Analytics Platform  |  Confidential",
              "size": 9, "color": "muted"}])
    _tblock(slide, SLIDE_W - 3.2, SLIDE_H - 0.38, 2.65, 0.3,
            [{"text": f"Slide {index} of {total}", "size": 9, "color": "muted",
              "align": PP_ALIGN.RIGHT}])


def _kv_cell(slide, x, y, w, h, label, value):
    _rect(slide, x, y, w, h, fill="light", rounded=True)
    _rect(slide, x, y, 0.07, h, fill="gold")
    _tblock(slide, x + 0.25, y + 0.04, w - 0.4, h,
            [
                {"text": _safe(label).upper(), "size": 8.5, "color": "muted"},
                {"text": _safe(value), "size": 12, "bold": True,
                 "color": "text", "space_before": 1},
            ])


def _kpi_card(slide, x, y, w, h, label, value, tone="gold", sub=""):
    _rect(slide, x, y, w, h, fill="light", rounded=True)
    _rect(slide, x, y, w, 0.07, fill=tone)
    paras = [
        {"text": _safe(label).upper(), "size": 9.5, "color": "muted"},
        {"text": _safe(value), "size": 20, "bold": True, "color": "text",
         "space_before": 2},
    ]
    if sub:
        paras.append({"text": _safe(sub), "size": 8.5, "color": "muted",
                      "space_before": 1})
    _tblock(slide, x + 0.22, y + 0.12, w - 0.44, h - 0.2, paras)


# ── Slides ──────────────────────────────────────────────────────────────────────


def _slide_title(prs, result, analysis, info, now):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill="navy")
    _rect(slide, 0, 0, SLIDE_W, 0.18, fill="gold")

    prob = analysis["prob"]
    label = result.get("label", "Likely to Stay")

    _tblock(slide, 0.6, 1.45, 12, 0.4,
            [{"text": "CUSTOMER CHURN ANALYTICS PLATFORM", "size": 15,
              "bold": True, "color": "gold_bright"}])
    _tblock(slide, 0.6, 2.0, 12.1, 1.6,
            [{"text": "Customer Churn Prediction Dashboard", "size": 42,
              "bold": True, "color": "white", "line_spacing": 1.02}])
    _tblock(slide, 0.6, 3.55, 12, 0.45,
            [{"text": "Executive Customer Retention Report", "size": 19,
              "color": "teal"}])
    _rect(slide, 0.62, 4.5, 2.6, 0.06, fill="gold")

    _tblock(slide, 0.6, 4.85, 12, 0.45,
            [{"text": f"{label}  -  {analysis['risk']}  -  "
                      f"{analysis['prob']:.1f}% churn probability",
              "size": 15, "bold": True, "color": "ink"}])

    _tblock(slide, 0.6, 5.55, 12, 0.4,
            [{"text": f"Generated: {now.strftime('%B %d, %Y')} at "
                      f"{now.strftime('%I:%M %p')}",
              "size": 12, "color": "sub"}])
    _tblock(slide, 0.6, 5.95, 12, 0.4,
            [{"text": "Prepared by Abhinav Agnihotri  -  "
                      "College Project Evaluation",
              "size": 12, "color": "sub"}])
    _tblock(slide, 0.6, 6.55, 12, 0.35,
            [{"text": f"Deployed model: {info.get('label', 'XGBoost')}",
              "size": 10.5, "color": "muted"}])


def _slide_executive_summary(prs, result, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Executive Summary",
            "Overall prediction, risk level, and probability")
    _footer(slide, 2, 6)

    prob = analysis["prob"]
    label = result.get("label", "Likely to Stay")
    risk = analysis["risk"]
    risk_color = RISK_COLOR.get(risk, "gold")

    card_w = (12.23 - 2 * 0.4) / 3.0
    cards = [
        ("Overall Prediction", label, VERDICT_COLOR.get(label, "gold")),
        ("Risk Level", risk, risk_color),
        ("Churn Probability", f"{prob:.1f}%",
         VERDICT_COLOR.get(label, "gold")),
    ]
    for i, (cap, value, tone) in enumerate(cards):
        x = 0.55 + i * (card_w + 0.4)
        _kpi_card(slide, x, 1.4, card_w, 1.15, cap, value, tone=tone)

    _tblock(slide, 0.55, 2.82, 8, 0.3,
            [{"text": "CHURN PROBABILITY", "size": 10, "bold": True,
              "color": "gold_dark"}])
    _tblock(slide, 10.2, 2.82, 2.6, 0.3,
            [{"text": f"{prob:.1f}%", "size": 16, "bold": True,
              "color": "text", "align": PP_ALIGN.RIGHT}])
    _rect(slide, 0.55, 3.25, 12.23, 0.4, fill="light", line="line")
    fill_w = 12.23 * min(max(prob, 0.0), 100.0) / 100.0
    _rect(slide, 0.55, 3.25, fill_w, 0.4, fill=risk_color)

    _tblock(slide, 0.55, 4.05, 7.7, 0.3,
            [{"text": "EXECUTIVE SUMMARY", "size": 12, "bold": True,
              "color": "gold_dark"}])
    if prob >= 70:
        stance = ("at immediate risk of churn and requires urgent "
                  "retention action")
    elif prob >= 40:
        stance = ("exposed to a meaningful churn risk that warrants "
                  "proactive outreach")
    else:
        stance = ("in a stable, low-risk position with strong retention "
                  "upside")
    paragraph = (
        f"This customer is {stance}. The deployed model estimates a "
        f"{prob:.1f}% churn probability ({risk}), with an estimated annual "
        f"revenue at risk of {_money(analysis['revenue_at_risk'])} and a "
        f"24-month customer-lifetime value estimate of "
        f"{_money(analysis['clv_estimate'])}."
    )
    if analysis.get("top_driver"):
        top = analysis["top_driver"]
        paragraph += (f" The strongest churn driver identified is "
                      f"{top['feature']} ({top['value']}).")
    _tblock(slide, 0.55, 4.4, 7.7, 2.5,
            [{"text": paragraph, "size": 12, "color": "body",
              "line_spacing": 1.25}])

    _rect(slide, 8.4, 4.05, 4.38, 2.6, fill="light", rounded=True)
    _tblock(slide, 8.65, 4.2, 3.9, 0.3,
            [{"text": "KEY HIGHLIGHTS", "size": 11, "bold": True,
              "color": "gold_dark"}])
    highlights = [
        f"Prediction: {label} ({prob:.1f}%)",
        f"Risk classification: {risk}  |  Priority: {analysis['priority']}",
        f"Segment: {analysis['segment']}",
        f"Estimated net value of plan: "
        f"{_money(analysis['net_benefit'])}  "
        f"({analysis['roi_pct']:.0f}% ROI)",
    ]
    _tblock(slide, 8.65, 4.6, 3.95, 2.0,
            [{"text": f"-  {h}", "size": 11, "color": "body",
              "space_after": 7, "line_spacing": 1.1} for h in highlights])


def _slide_customer_details(prs, inputs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Customer Details",
            "The customer profile used to generate this prediction")
    _footer(slide, 3, 6)

    pairs = [
        (prediction.FEATURE_LABELS.get(feat, feat),
         prediction.display_value(feat, inputs[feat]))
        for feat in prediction.FEATURE_NAMES
    ]
    col_w = 5.98
    x0 = 0.55
    x1 = 0.55 + col_w + 0.27
    row_h = 0.5
    y0 = 1.35
    for i, (label, value) in enumerate(pairs):
        col = i // 10
        row = i % 10
        x = x0 if col == 0 else x1
        y = y0 + row * row_h
        _kv_cell(slide, x, y, col_w, row_h - 0.06, label, value)


def _slide_kpi_summary(prs, result, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "KPI Summary",
            "Key modeled financial and account metrics")
    _footer(slide, 4, 6)

    kpis = [
        ("Churn Probability", f"{analysis['prob']:.1f}%", "red"),
        ("Revenue at Risk (12-mo)", _money(analysis["revenue_at_risk"]), "red"),
        ("CLV Estimate (24-mo)", _money(analysis["clv_estimate"]), "blue"),
        ("Retention Investment", _money(analysis["retention_cost"]), "gold"),
        ("Modeled Benefit (65%)", _money(analysis["potential_savings"]), "green"),
        ("Net Expected Value", _money(analysis["net_benefit"]), "green"),
        ("Estimated ROI", f"{analysis['roi_pct']:.0f}%", "emerald"),
        ("Priority", analysis["priority"], "purple"),
        ("Confidence", analysis["confidence"], "blue"),
    ]
    card_w = 3.81
    card_h = 0.98
    gap_x = 0.4
    gap_y = 0.28
    x0 = 0.55
    y0 = 1.35
    for i, (label, value, tone) in enumerate(kpis):
        row = i // 3
        col = i % 3
        x = x0 + col * (card_w + gap_x)
        y = y0 + row * (card_h + gap_y)
        _kpi_card(slide, x, y, card_w, card_h, label, value, tone=tone)

    insights_y = y0 + 3 * (card_h + gap_y) + 0.18
    _rect(slide, 0.55, insights_y, 12.23, 1.9, fill="light", rounded=True)
    _tblock(slide, 0.85, insights_y + 0.14, 11.5, 0.3,
            [{"text": "BUSINESS INSIGHTS", "size": 12, "bold": True,
              "color": "gold_dark"}])
    insights = [
        f"The account carries a {analysis['prob']:.0f}% churn probability "
        f"({analysis['risk']}) and warrants a {analysis['priority']} response.",
        f"An investment of {_money(analysis['retention_cost'])} is expected "
        f"to protect {_money(analysis['revenue_at_risk'])} of annual revenue "
        f"({_money(analysis['net_benefit'])} net value, "
        f"{analysis['roi_pct']:.0f}% ROI).",
        f"The modeled CLV over 24 months is "
        f"{_money(analysis['clv_estimate'])}.",
    ]
    _tblock(slide, 0.85, insights_y + 0.5, 11.7, 1.3,
            [{"text": f"-  {t}", "size": 11.5, "color": "body",
              "space_after": 6, "line_spacing": 1.15} for t in insights])


def _slide_recommendations(prs, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Business Recommendations",
            "Prioritized action items for this account")
    _footer(slide, 5, 6)

    actions = analysis.get("actions") or []
    if not actions:
        _tblock(slide, 0.55, 1.8, 12, 0.5,
                [{"text": "No specific recommendations are available for "
                          "this account at this time.", "size": 13,
                  "color": "body"}])
    for i, action in enumerate(actions[:6]):
        y = 1.45 + i * 0.82
        icon = action.get("icon", "")
        num = f"{icon} {i + 1}" if icon else str(i + 1)
        _rect(slide, 0.55, y, 12.23, 0.7, fill="light", rounded=True)
        _rect(slide, 0.55, y, 0.09, 0.7, fill="gold")
        _tblock(slide, 0.85, y + 0.06, 1.1, 0.55,
                [{"text": num, "size": 15, "bold": True, "color": "gold_dark",
                  "align": PP_ALIGN.CENTER}])
        title = action.get("title", "Action item")
        _tblock(slide, 2.05, y + 0.05, 10.5, 0.32,
                [{"text": title, "size": 13.5, "bold": True, "color": "text"}])
        detail = []
        if action.get("reason"):
            detail.append(f"Business reason: {action['reason']}")
        if action.get("impact"):
            detail.append(f"Impact: {action['impact']}")
        if action.get("cost"):
            detail.append(f"Cost: {action['cost']}")
        _tblock(slide, 2.05, y + 0.36, 10.5, 0.3,
                [{"text": "   ".join(detail), "size": 9.5, "color": "body"}])

    campaign = analysis.get("campaign") or {}
    if campaign.get("name"):
        _rect(slide, 0.55, 6.35, 12.23, 0.55, fill="navy", rounded=True)
        _tblock(slide, 0.85, 6.4, 11.7, 0.42,
                [{"text": f"RECOMMENDED CAMPAIGN:  "
                          f"{_safe(campaign.get('name'))}",
                  "size": 12.5, "bold": True, "color": "gold_bright"}])


def _slide_conclusion(prs, result, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Conclusion",
            "Summary judgment and recommended next steps")
    _footer(slide, 6, 6)

    prob = analysis["prob"]
    if prob >= 70:
        judgment = "immediate retention intervention"
    elif prob >= 40:
        judgment = "a proactive retention campaign"
    else:
        judgment = "continued relationship building and growth"

    _rect(slide, 0.55, 1.45, 12.23, 2.3, fill="light", rounded=True)
    _rect(slide, 0.55, 1.45, 12.23, 0.07, fill="gold")
    _tblock(slide, 0.85, 1.7, 11.6, 0.35,
            [{"text": "CONCLUSION", "size": 13, "bold": True,
              "color": "gold_dark"}])
    conclusion = (
        f"This account requires {judgment}. With a {prob:.1f}% churn "
        f"probability and an estimated annual revenue exposure of "
        f"{_money(analysis['revenue_at_risk'])}, acting on the recommended "
        f"plan - a {_money(analysis['retention_cost'])} investment for a "
        f"modeled benefit of {_money(analysis['potential_savings'])} and an "
        f"ROI of {analysis['roi_pct']:.0f}% - offers a compelling business "
        f"case."
    )
    _tblock(slide, 0.85, 2.1, 11.6, 1.5,
            [{"text": conclusion, "size": 12.5, "color": "body",
              "line_spacing": 1.25}])

    _tblock(slide, 0.55, 4.15, 12, 0.35,
            [{"text": "NEXT STEPS", "size": 13, "bold": True,
              "color": "gold_dark"}])
    next_steps = [
        f"Execute the recommended plan: "
        f"{_safe((analysis.get('campaign') or {}).get('name', 'Retention follow-up'))}.",
        "Review this customer at the next weekly retention desk meeting.",
        "Track the account for 30 days and re-run the model on refreshed inputs.",
        "Escalate if the churn probability moves above the priority threshold.",
    ]
    _tblock(slide, 0.55, 4.6, 12.2, 1.3,
            [{"text": f"{i + 1}.  {s}", "size": 11.5, "color": "body",
              "space_after": 5} for i, s in enumerate(next_steps)])

    if analysis.get("manager_notes"):
        _rect(slide, 0.55, 6.2, 12.23, 0.75, fill="navy", rounded=True)
        _tblock(slide, 0.85, 6.28, 11.6, 0.6,
                [{"text": f"ACCOUNT MANAGER NOTE:  "
                          f"{_safe(analysis['manager_notes'])}",
                  "size": 10.5, "color": "ink", "line_spacing": 1.15}])


# ── Public entry point ─────────────────────────────────────────────────────────


def build_presentation(result: dict, inputs: dict,
                       a: dict | None = None) -> bytes:
    """Build the 6-slide executive deck and return it as bytes.

    `result` is the prediction dict from `prediction.predict()` (plus
    `factors`, `model_alias`, and `recommendations`). `inputs` is the raw
    customer form values. `a` is the optional recommendation-engine analysis;
    missing keys fall back to values derived from the prediction.
    """
    analysis = report.resolve_analysis(result, inputs, a)
    info = prediction.model_info(
        result.get("model_alias", prediction.DEFAULT_MODEL)
    )
    now = datetime.now()

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    _slide_title(prs, result, analysis, info, now)
    _slide_executive_summary(prs, result, analysis)
    _slide_customer_details(prs, inputs)
    _slide_kpi_summary(prs, result, analysis)
    _slide_recommendations(prs, analysis)
    _slide_conclusion(prs, result, analysis)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
